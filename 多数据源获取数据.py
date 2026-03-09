# Dify 依赖管理: markitdown-no-magika, httpx, trafilatura, beautifulsoup4, lxml, PyMuPDF, pdfplumber, python-pptx, python-docx, openpyxl, xlrd, Pillow
import asyncio
import httpx
import re
import os
import json
import time
import traceback
from typing import Any, Dict, List, Literal, Optional
from abc import ABC, abstractmethod
from io import BytesIO
from urllib.parse import urljoin
import tempfile
import base64
import hashlib

try:
    from markitdown import MarkItDown  # pip install markitdown-no-magika (无 onnxruntime 依赖)
except ImportError as e:
    print(f"‼️ MarkItDown Import Error: {e}")
    MarkItDown = None
except Exception as e:
    print(f"‼️ MarkItDown Unexpected Error: {e}")
    MarkItDown = None

# --- 核心依赖 ---
# trafilatura 用于从HTML提取主要内容
import trafilatura
from trafilatura.settings import use_config

# PyPDF2 用于解析PDF
import pdfplumber
import fitz
# BeautifulSoup 用于辅助解析HTML（例如提取视频）
from bs4 import BeautifulSoup
import csv
import zipfile
import xml.etree.ElementTree as ET

# --- 文件解析增强依赖 (可选, 缺失时自动降级到 MarkItDown) ---
try:
    from docx import Document as DocxDocument
except ImportError:
    DocxDocument = None

try:
    from pptx import Presentation as PptxPresentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE as PptxShapeType
except ImportError:
    PptxPresentation = None
    PptxShapeType = None

try:
    from openpyxl import load_workbook as openpyxl_load_workbook
except ImportError:
    openpyxl_load_workbook = None

try:
    import xlrd
except ImportError:
    xlrd = None

try:
    from PIL import Image as PILImage
except ImportError:
    PILImage = None

# ==============================================================================
# ====================== DIFY 本地调试辅助模块 =========================
# ==============================================================================
import pprint

# --- 本地调试开关 ---
# 在你的 IDE 中进行测试时，将此值设为 True。
# 当你准备将代码复制到 Dify 平台时，请将其改回 False，或直接删除此调试模块。
IS_LOCAL_DEBUG = True


def _dify_debug_return(data: Dict[str, Any], label: str = "Final Return") -> Dict[str, Any]:
    """
    一个用于在 Dify 代码节点中进行本地调试的包装函数。

    当 IS_LOCAL_DEBUG 为 True 时，它会漂亮地打印出最终要返回的数据，
    然后原封不动地返回该数据，以便 Dify 平台能正确接收。

    Args:
        data (Dict[str, Any]): 准备从 Dify 节点返回的数据。
        label (str, optional): 一个标签，用于在控制台输出中标识来源。默认为 "Final Return"。

    Returns:
        Dict[str, Any]: 传入的原始数据。
    """
    if IS_LOCAL_DEBUG:
        # 打印一个清晰的分隔符和标签，方便在终端中识别
        print("\n" + "=" * 40 + f" DIFY DEBUG OUTPUT [{label}] " + "=" * 40)

        # 使用 pprint 模块进行美化输出，对复杂的嵌套字典特别友好
        pprint.pprint(data, indent=2, width=120)

        # 打印结束分隔符
        print("=" * 105 + "\n")

    # 无论是否打印，都必须原封不动地返回原始数据
    return data


def _parse_input_data(raw_input: Any) -> Dict[str, Any]:
    """
    健壮地解析上一个节点的输出，能同时处理带 "datas" 包装和不带包装的两种结构。
    并分离出不同来源的网页搜索URL、视频URL、招聘查询参数和企业名称，同时保留元数据。
    """
    print(
        f"============== 步骤 1: 接收到原始输入 ==============\nTYPE: {type(raw_input)}\nVALUE: {raw_input}\n=======================================================")
    if isinstance(raw_input, str):
        if not raw_input.strip(): return {"web_url_info_list": [], "video_url_info_list": [], "career_payload": {},
                                          "enterprise_names": []}
        try:
            data = json.loads(raw_input)
        except json.JSONDecodeError as e:
            raise ValueError(f"无法将输入字符串解析为JSON: {e}")
    elif isinstance(raw_input, dict):
        data = raw_input
    else:
        raise TypeError(f"期望的输入类型是 str 或 dict, 但收到了 {type(raw_input).__name__}")

    if "datas" in data and isinstance(data["datas"], dict):
        print("  [解析器] 检测到 'datas' 包装层，将使用其内部数据。")
        datas_obj = data["datas"]
    else:
        print("  [解析器] 未检测到 'datas' 包装层，将直接使用顶层数据。")
        datas_obj = data

    if not isinstance(datas_obj, dict): datas_obj = {}

    # Extract run_mode from input data if available
    extracted_run_mode = datas_obj.get("run_mode")
    print(f"  [解析器] 从输入中提取到的 run_mode: {extracted_run_mode}")

    # --- 1. 提取 Legacy Mode 数据 ---
    comprehensive_data = datas_obj.get("comprehensive_data", [])
    tianyan_data = datas_obj.get("tianyan_check_data", [])

    # --- 2. 提取 External Mode 数据 ---
    general_web_data = datas_obj.get("general_web_data", [])
    institution_source_data = datas_obj.get("institution_source_data", [])
    # 兼容 Tuoyu 模式下 web_query 可能被分配到的位置，或者未来直接传递 tuoyu_web_data
    # 目前上游逻辑是将 tuoyu_web_queries 分配给了 general_web_data 或 institution_source_data
    # 所以这里不需要额外提取 tuoyu_web_data，除非上游结构改变

    # --- 3. 提取 Shared 数据 ---
    career_data = datas_obj.get("career_data", {})

    web_url_info_list = []
    video_url_info_list = []

    def _extract_urls(source_data, origin_key):
        """Helper to extract URLs from a list of query results."""
        if not isinstance(source_data, list): return
        for query_result in source_data:
            if not isinstance(query_result, dict): continue
            query_text = query_result.get("query", "")

            # Iterate through all keys to find result lists (e.g., 'web_results', 'video_results', 'policy_regional_results')
            for key, result_list in query_result.items():
                if not key.endswith("_results") or not isinstance(result_list, list):
                    continue

                for res in result_list:
                    if not isinstance(res, dict): continue

                    # Identify provider
                    provider = next(
                        (k[:-4] for k in res if k.endswith('_url') and '_embed_' not in k and '_thumbnail_' not in k),
                        None)
                    if not provider:
                        provider = next((k.split('_')[0] for k in res if '_url' in k), None)
                        if not provider: continue

                    result_type = res.get(f"{provider}_type")

                    # Common info
                    info = {
                        "url": res.get(f"{provider}_url"),
                        "title": res.get(f"{provider}_title", "Untitled"),
                        "source": res.get(f"{provider}_source"),
                        "snippet": res.get(f"{provider}_snippet"),
                        "provider": provider,
                        "query": query_text,
                        "origin_key": origin_key  # Tag the source origin
                    }

                    if not info["url"]: continue

                    if result_type == 'video':
                        info.update({
                            "video_id": res.get(f"{provider}_video_id"),
                            "embed_url": res.get(f"{provider}_embed_url"),
                            "thumbnail_url": res.get(f"{provider}_thumbnail_url"),
                        })
                        video_url_info_list.append(info)
                    else:
                        web_url_info_list.append(info)

    # Process all sources
    _extract_urls(comprehensive_data, "comprehensive_data")
    _extract_urls(general_web_data, "general_web_data")
    _extract_urls(institution_source_data, "institution_source_data")

    career_payload = career_data if isinstance(career_data, dict) else {}
    enterprise_names: List[str] = []
    if isinstance(tianyan_data, str) and tianyan_data.strip():
        print("  [解析器] 检测到 tianyan_check_data 为字符串，将处理单个企业。")
        enterprise_names.append(tianyan_data.strip())
    elif isinstance(tianyan_data, list):
        print(f"  [解析器] 检测到 tianyan_check_data 为列表，将处理 {len(tianyan_data)} 个企业。")
        enterprise_names = [str(name).strip() for name in tianyan_data if isinstance(name, str) and str(name).strip()]

    # Detect Mode
    mode = "legacy"
    if extracted_run_mode and extracted_run_mode.lower() == "tuoyu":
        mode = "external"
    elif extracted_run_mode and extracted_run_mode.lower() == "x-pilot":
        mode = "legacy"
    else:
        if general_web_data or institution_source_data:
            mode = "external"

    parsed_result = {
        "mode": mode,
        "web_url_info_list": web_url_info_list,
        "video_url_info_list": video_url_info_list,
        "career_payload": career_payload,
        "enterprise_names": enterprise_names
    }

    print(
        f"============== 步骤 2: 输入解析完毕 ==============\n模式: {mode}\n网页URL数量: {len(web_url_info_list)}\n视频URL数量: {len(video_url_info_list)}\n招聘负载: {career_payload}\n企业名称列表: {enterprise_names} (共 {len(enterprise_names)} 个)\n=======================================================")

    return parsed_result


# --- 1. 输入解析模块 ---
# 【已修复】替换这个函数
# def _parse_input_data(raw_input: Any) -> Dict[str, Any]:
#     """
#     健壮地解析上一个节点的输出，能同时处理带 "datas" 包装和不带包装的两种结构。
#     """
#     print(
#         f"============== 步骤 1: 接收到原始输入 ==============\nTYPE: {type(raw_input)}\nVALUE: {raw_input}\n=======================================================")
#     if isinstance(raw_input, str):
#         if not raw_input.strip(): return {"url_list": [], "career_payload": {}, "enterprise_name": ""}
#         try:
#             data = json.loads(raw_input)
#         except json.JSONDecodeError as e:
#             raise ValueError(f"无法将输入字符串解析为JSON: {e}")
#     elif isinstance(raw_input, dict):
#         data = raw_input
#     else:
#         raise TypeError(f"期望的输入类型是 str 或 dict, 但收到了 {type(raw_input).__name__}")
#
#     # --- 核心修复逻辑 ---
#     # 检查顶层是否有 "datas" 键，如果没有，就认为当前整个对象就是我们要的数据体。
#     if "datas" in data and isinstance(data["datas"], dict):
#         print("  [解析器] 检测到 'datas' 包装层，将使用其内部数据。")
#         datas_obj = data["datas"]
#     else:
#         print("  [解析器] 未检测到 'datas' 包装层，将直接使用顶层数据。")
#         datas_obj = data
#     # --- 修复结束 ---
#
#     if not isinstance(datas_obj, dict): datas_obj = {}
#
#     comprehensive_data = datas_obj.get("comprehensive_data", [])
#     career_data = datas_obj.get("career_data", {})
#     tianyan_data = datas_obj.get("tianyan_check_data", "")
#
#     url_list = []
#     if isinstance(comprehensive_data, list):
#         for query_result in comprehensive_data:
#             if not isinstance(query_result, dict): continue
#             for res_list_key in ["web_results", "video_results"]:
#                 for res in query_result.get(res_list_key, []):
#                     if not isinstance(res, dict): continue
#                     url, title, provider = None, None, None
#                     for key, value in res.items():
#                         if key.endswith("_url"):
#                             url, provider = value, key.split('_url')[0]
#                             title = res.get(f"{provider}_title", "Untitled")
#                             break
#                     if url and provider:
#                         url_list.append({"url": url, "title": title, "provider": provider})
#     career_payload = career_data if isinstance(career_data, dict) else {}
#     enterprise_name = tianyan_data if isinstance(tianyan_data, str) else ""
#     parsed_result = {"url_list": url_list, "career_payload": career_payload, "enterprise_name": enterprise_name.strip()}
#
#     print(
#         f"============== 步骤 2: 输入解析完毕 ==============\nURL 数量: {len(url_list)}\n招聘负载: {career_payload}\n企业名称: '{enterprise_name.strip()}'\n=======================================================")
#
#     return parsed_result


# def _parse_input_data(raw_input: Any) -> Dict[str, Any]:
#     """
#     健壮地解析上一个节点的输出，分离出web搜索URL和招聘查询参数。
#     """
#     if isinstance(raw_input, str):
#         if not raw_input.strip(): return {"url_list": [], "career_payload": {}}
#         try:
#             data = json.loads(raw_input)
#         except json.JSONDecodeError as e:
#             raise ValueError(f"无法将输入字符串解析为JSON: {e}")
#     elif isinstance(raw_input, dict):
#         data = raw_input
#     else:
#         raise TypeError(f"期望的输入类型是 str 或 dict, 但收到了 {type(raw_input).__name__}")
#     # 安全地深入到 'datas' 结构
#     datas_obj = data.get("datas", {})
#     if not isinstance(datas_obj, dict): datas_obj = {}
#     comprehensive_data = datas_obj.get("comprehensive_data", [])
#     career_data = datas_obj.get("career_data", {})
#     tianyan_data = datas_obj.get("tianyan_check_data", "")
#     # 1. 提取URL列表
#     url_list = []
#     if isinstance(comprehensive_data, list):
#         for query_result in comprehensive_data:
#             if not isinstance(query_result, dict): continue
#             for res_list_key in ["web_results", "video_results"]:
#                 for res in query_result.get(res_list_key, []):
#                     if not isinstance(res, dict): continue
#                     url, title, provider = None, None, None
#                     for key, value in res.items():
#                         if key.endswith("_url"):
#                             url = value
#                             provider = key.split('_url')[0]
#                             title = res.get(f"{provider}_title", "Untitled")
#                             break
#                     if url and provider:
#                         url_list.append({"url": url, "title": title, "provider": provider})
#     # 2. 提取职业查询负载
#     career_payload = career_data if isinstance(career_data, dict) else {}

#     # 3. 提取企业名称
#     enterprise_name = tianyan_data if isinstance(tianyan_data, str) else ""
#     return {"url_list": url_list, "career_payload": career_payload, "enterprise_name": enterprise_name.strip()}

# --- 2. 抽象与实现分离：内容抓取器 ---

# ==============================================================================
# ============== 嵌入图片上传服务 (EmbeddedImageUploader) ==============
# ==============================================================================

class EmbeddedImageUploader:
    """
    将文档内嵌图片上传到服务器，返回可访问的 URL。
    上传接口: POST /file/uploads (multipart/form-data, 字段名 files)
    响应格式: {"status": true, "data": [{"originalname": "x.png", "url": "https://..."}]}
    """
    UPLOAD_URL = os.environ.get("IMAGE_UPLOAD_URL", "https://server.x-pilot.cn/file/uploads")
    MAX_BATCH = 10
    SUPPORTED_EXTS = {'.png', '.jpg', '.jpeg', '.gif', '.bmp', '.webp', '.tiff', '.tif'}
    MIME_MAP = {
        '.png': 'image/png', '.jpg': 'image/jpeg', '.jpeg': 'image/jpeg',
        '.gif': 'image/gif', '.bmp': 'image/bmp', '.webp': 'image/webp',
        '.tiff': 'image/tiff', '.tif': 'image/tiff', '.svg': 'image/svg+xml',
    }

    @classmethod
    def upload_images(cls, images: List[tuple]) -> Dict[str, str]:
        """
        批量上传图片，返回 {原始文件名: 可访问URL} 映射。
        images: [(filename, binary_data, mime_type), ...]
        """
        if not images:
            return {}
        url_map: Dict[str, str] = {}
        for i in range(0, len(images), cls.MAX_BATCH):
            batch = images[i:i + cls.MAX_BATCH]
            files_payload = [('files', (fname, data, mime)) for fname, data, mime in batch]
            try:
                with httpx.Client(timeout=60, verify=False) as client:
                    resp = client.post(cls.UPLOAD_URL, files=files_payload)
                    resp.raise_for_status()
                    resp_data = resp.json()
                    if resp_data.get("status") and resp_data.get("data"):
                        for item in resp_data["data"]:
                            orig = item.get("originalname", "")
                            url = item.get("url", "")
                            if orig and url:
                                url_map[orig] = url
                                print(f"    📤 已上传: {orig} -> {url}")
            except Exception as e:
                print(f"⚠️ 图片批次上传失败: {e}")
        return url_map

    @classmethod
    def extract_from_zip(cls, data: bytes, media_prefix: str, min_size: int = 5120) -> List[tuple]:
        """从 ZIP 格式文档 (docx/pptx) 中提取 media 目录下的图片。"""
        images = []
        try:
            with zipfile.ZipFile(BytesIO(data)) as zf:
                for name in zf.namelist():
                    if not name.startswith(media_prefix):
                        continue
                    basename = os.path.basename(name)
                    ext_lower = os.path.splitext(basename)[1].lower()
                    if ext_lower not in cls.SUPPORTED_EXTS:
                        continue
                    img_data = zf.read(name)
                    if len(img_data) < min_size:
                        continue
                    mime = cls.MIME_MAP.get(ext_lower, 'image/png')
                    images.append((basename, img_data, mime))
        except Exception as e:
            print(f"⚠️ ZIP 图片提取失败: {e}")
        return images

    @classmethod
    def extract_from_pdf(cls, data: bytes, max_pages: int = 50, min_size: int = 5120, min_dim: int = 50) -> List[tuple]:
        """从 PDF 中提取嵌入图片 (fitz)，按页面顺序返回。"""
        images = []
        try:
            with fitz.open(stream=data, filetype="pdf") as doc:
                img_idx = 0
                for pi in range(min(len(doc), max_pages)):
                    page = doc.load_page(pi)
                    page_dict = page.get_text("dict", sort=True)
                    for block in page_dict.get("blocks", []):
                        if block["type"] != 1:
                            continue
                        bbox = block.get("bbox", [0, 0, 0, 0])
                        w, h = bbox[2] - bbox[0], bbox[3] - bbox[1]
                        if w < min_dim or h < min_dim:
                            continue
                        img_bytes = block.get("image", b"")
                        if len(img_bytes) < min_size:
                            continue
                        img_idx += 1
                        ext = "png"
                        if img_bytes[:3] == b'\xff\xd8\xff':
                            ext = "jpg"
                        fname = f"pdf_image_{img_idx}.{ext}"
                        images.append((fname, img_bytes, f"image/{ext}"))
        except Exception as e:
            print(f"⚠️ PDF 图片提取失败: {e}")
        return images


# ==============================================================================
# ================ 数据清洗管道 (DataCleaningPipeline) ================
# ==============================================================================

class DataCleaningPipeline:
    """LLM 友好的多阶段数据清洗管道，所有输出均为干净的 Markdown 字符串。"""

    _NOISY_PATTERNS = [re.compile(p, re.IGNORECASE) for p in [
        r'^[\-=*#_]{3,}$',
        r'.*\.(html|shtml|htm|php)\s*$',
        r'.{0,50}(搜狐|网易|腾讯|新浪|登录|注册|版权所有|版权声明).{0,50}$',
        r'\[\d+\]|\[下一页\]|\[上一页\]',
        r'\[(编辑|查看历史|讨论|阅读|来源|原标题)\]',
        r'^\*+\s*\[.*?\]\(.*?\)',
        r'^\s*(分享到|扫描二维码|返回搜狐|查看更多|责任编辑|记者|通讯员)',
        r'^\s*([京公网安备京网文京ICP备]|互联网新闻信息服务许可证|信息网络传播视听节目许可证)',
    ]]
    _IMG_PATTERN = re.compile(r'(!\[(.*?)\]\((.*?)\))')
    _LINK_PATTERN = re.compile(r'\[.*?\]\(.*?\)')
    _EDITOR_PATTERN = re.compile(r'(\(|\[)\s*责任编辑：.*?\s*(\)|\])')
    _PAGE_NUM_PATTERN = re.compile(
        r'^\s*[-—]\s*\d+\s*[-—]\s*$|'
        r'^\s*第\s*\d+\s*页\s*(共\s*\d+\s*页)?\s*$|'
        r'^\s*Page\s+\d+\s*(of\s+\d+)?\s*$',
        re.IGNORECASE
    )
    _REPEATED_LINE_THRESHOLD = 3

    def __init__(self, max_content_length: int = 80000):
        self.max_content_length = max_content_length

    @classmethod
    def _is_noisy_line(cls, line: str) -> bool:
        stripped = line.strip()
        if not stripped:
            return True
        for pat in cls._NOISY_PATTERNS:
            if pat.search(stripped):
                return True
        links = cls._LINK_PATTERN.findall(stripped)
        if len(links) > 2 and len(stripped) / (len(links) + 1) < 30:
            return True
        return False

    @staticmethod
    def _normalize_whitespace(text: str) -> str:
        lines = text.splitlines()
        out, prev_empty = [], False
        for line in lines:
            stripped = line.strip()
            if not stripped:
                if not prev_empty:
                    out.append("")
                prev_empty = True
            else:
                out.append(stripped)
                prev_empty = False
        return "\n".join(out).strip()

    @classmethod
    def _remove_repeated_headers_footers(cls, text: str) -> str:
        lines = text.splitlines()
        if len(lines) < 20:
            return text
        line_counts: Dict[str, int] = {}
        for line in lines:
            s = line.strip()
            if s and len(s) < 100:
                line_counts[s] = line_counts.get(s, 0) + 1
        repeated = {s for s, c in line_counts.items() if c >= cls._REPEATED_LINE_THRESHOLD}
        if not repeated:
            return text
        return "\n".join(l for l in lines if l.strip() not in repeated)

    def _truncate(self, text: str, label: str = "内容") -> str:
        if len(text) > self.max_content_length:
            return text[:self.max_content_length] + f"\n\n...[{label}过长，已截断至 {self.max_content_length} 字符]"
        return text

    def clean_document(self, text: str) -> str:
        if not text:
            return ""
        text = self._remove_repeated_headers_footers(text)
        lines = text.splitlines()
        cleaned = []
        for line in lines:
            if self._PAGE_NUM_PATTERN.search(line.strip()):
                continue
            if self._is_noisy_line(line):
                continue
            line = self._EDITOR_PATTERN.sub('', line).strip()
            if line:
                cleaned.append(line)
        result = self._normalize_whitespace("\n".join(cleaned))
        return self._truncate(result, "文档内容")

    def clean_html(self, text: str) -> str:
        if not text:
            return ""
        lines = text.splitlines()
        cleaned = []
        for line in lines:
            if self._is_noisy_line(line):
                continue
            line = self._EDITOR_PATTERN.sub('', line).strip()
            if line:
                cleaned.append(line)
        result = self._normalize_whitespace("\n".join(cleaned))
        return self._truncate(result, "网页内容")

    def clean_table(self, text: str) -> str:
        if not text:
            return ""
        text = self._normalize_whitespace(text)
        return self._truncate(text, "表格内容")

    def clean_text(self, text: str) -> str:
        if not text:
            return ""
        text = self._normalize_whitespace(text)
        return self._truncate(text)

    async def validate_image_urls(self, md_text: str, client: httpx.AsyncClient) -> str:
        MAX_TO_VALIDATE = 25
        matches = list(self._IMG_PATTERN.finditer(md_text))
        if not matches:
            return md_text
        urls_all = {m.group(3).strip() for m in matches}
        urls_to_check = set(list(urls_all)[:MAX_TO_VALIDATE])

        async def _check(u):
            if not u or not u.startswith(('http://', 'https://')):
                return u, False
            try:
                resp = await client.head(u, timeout=5, follow_redirects=True)
                ct = resp.headers.get('content-type', '').lower()
                return u, resp.is_success and 'image' in ct
            except Exception:
                return u, False

        results = await asyncio.gather(*[_check(u) for u in urls_to_check], return_exceptions=True)
        valid = set()
        for r in results:
            if isinstance(r, tuple) and r[1]:
                valid.add(r[0])
        valid.update(urls_all - urls_to_check)

        def _replacer(m):
            return m.group(0) if m.group(3).strip() in valid else ""

        return self._IMG_PATTERN.sub(_replacer, md_text)


# ==============================================================================
# ============ 统一文件解析服务 (DocumentParserService) ============
# ==============================================================================

class DocumentParserService:
    """
    统一文件解析服务，替代原 ResourceParser。
    支持: pdf, docx, doc, pptx, ppt, xlsx, xls, csv,
          html, htm, json, xml, txt, md,
          jpg, jpeg, png, gif, webp, bmp
    所有输出均为 LLM 友好的 Markdown 字符串。
    """

    PDF_MAX_PAGES = 50
    MAX_TABLE_ROWS = 500
    MAX_TEXT_CHARS = 100000
    MAX_JSON_CHARS = 50000
    MIN_IMG_BYTES = 5 * 1024
    MIN_IMG_DIM = 50

    def __init__(self):
        self.cleaner = DataCleaningPipeline()
        self._markitdown = None
        if MarkItDown:
            try:
                self._markitdown = MarkItDown()
            except Exception as e:
                print(f"⚠️ MarkItDown 初始化失败: {e}")

    # ── MarkItDown 通用转换 ──────────────────────────────────
    def _markitdown_convert(self, data: bytes, suffix: str) -> str:
        if not self._markitdown:
            return ""
        if not suffix.startswith("."):
            suffix = f".{suffix}"
        tmp_path = None
        try:
            with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as f:
                f.write(data)
                tmp_path = f.name
            result = self._markitdown.convert(tmp_path)
            return result.text_content if result and result.text_content else ""
        except Exception as e:
            print(f"⚠️ MarkItDown ({suffix}) 失败: {e}")
            return ""
        finally:
            if tmp_path and os.path.exists(tmp_path):
                try:
                    os.remove(tmp_path)
                except OSError:
                    pass

    # ── PDF ──────────────────────────────────────────────────
    @staticmethod
    def _bbox_overlap(bbox_a, bbox_b, tolerance=2.0) -> bool:
        """判断两个 bbox 是否在 Y 轴方向上有足够重叠 (用于去重表格区域内的散碎文本)"""
        ax0, ay0, ax1, ay1 = bbox_a
        bx0, by0, bx1, by1 = bbox_b
        if ax1 < bx0 + tolerance or bx1 < ax0 + tolerance:
            return False
        if ay1 < by0 + tolerance or by1 < ay0 + tolerance:
            return False
        overlap_x = min(ax1, bx1) - max(ax0, bx0)
        width_a = ax1 - ax0
        return width_a > 0 and (overlap_x / width_a) > 0.5

    def _parse_pdf(self, data: bytes, source_url: str = "") -> str:
        parts = []
        img_count = 0
        table_bboxes_per_page: Dict[int, list] = {}

        plumber_tables_per_page: Dict[int, list] = {}
        try:
            with pdfplumber.open(BytesIO(data)) as plumber_pdf:
                page_limit = min(len(plumber_pdf.pages), self.PDF_MAX_PAGES)
                for pi in range(page_limit):
                    pp = plumber_pdf.pages[pi]
                    tables = pp.find_tables()
                    if not tables:
                        continue
                    page_tables = []
                    page_bboxes = []
                    for tbl in tables:
                        rows = tbl.extract()
                        if not rows:
                            continue
                        cleaned = []
                        for row in rows:
                            cleaned.append([(c or "").strip() for c in row])
                        if any(any(cell for cell in r) for r in cleaned):
                            page_tables.append((tbl.bbox[1], cleaned))
                            page_bboxes.append(tbl.bbox)
                    if page_tables:
                        plumber_tables_per_page[pi] = page_tables
                        table_bboxes_per_page[pi] = page_bboxes
        except Exception as e:
            print(f"  ⚠️ pdfplumber 表格提取异常 (不影响正文): {e}")

        try:
            with fitz.open(stream=data, filetype="pdf") as doc:
                total = len(doc)
                limit = min(total, self.PDF_MAX_PAGES)
                if total > self.PDF_MAX_PAGES:
                    print(f"  📄 PDF 共 {total} 页，只处理前 {limit} 页")

                for pi in range(limit):
                    page = doc.load_page(pi)
                    page_dict = page.get_text("dict", sort=True)
                    tbl_bboxes = table_bboxes_per_page.get(pi, [])
                    elements = []

                    for block in page_dict.get("blocks", []):
                        b_bbox = block.get("bbox", [0, 0, 0, 0])
                        y0 = b_bbox[1]

                        if block["type"] == 0:
                            if tbl_bboxes and any(self._bbox_overlap(b_bbox, tb) for tb in tbl_bboxes):
                                continue
                            lines_text = []
                            for ln in block.get("lines", []):
                                span_txt = "".join(s.get("text", "") for s in ln.get("spans", []))
                                if span_txt.strip():
                                    lines_text.append(span_txt.strip())
                            if lines_text:
                                elements.append((y0, "\n".join(lines_text)))

                        elif block["type"] == 1:
                            w, h = b_bbox[2] - b_bbox[0], b_bbox[3] - b_bbox[1]
                            if w < self.MIN_IMG_DIM or h < self.MIN_IMG_DIM:
                                continue
                            img_bytes = block.get("image", b"")
                            if len(img_bytes) < self.MIN_IMG_BYTES:
                                continue
                            img_count += 1
                            elements.append(
                                (y0, f"![图片{img_count} (第{pi + 1}页, {int(w)}x{int(h)})](pdf_image_{img_count})"))

                    for tbl_y0, tbl_rows in plumber_tables_per_page.get(pi, []):
                        elements.append((tbl_y0, self._rows_to_md_table(tbl_rows)))

                    elements.sort(key=lambda x: x[0])
                    page_content = "\n\n".join(e[1] for e in elements)
                    if page_content.strip():
                        if limit > 1:
                            parts.append(f"<!-- 第 {pi + 1} 页 -->\n\n{page_content}")
                        else:
                            parts.append(page_content)

                if total > self.PDF_MAX_PAGES:
                    parts.append(f"\n\n> PDF 共 {total} 页，已处理前 {limit} 页")

            result = "\n\n".join(parts).strip()
            if result:
                result = self._upload_embedded_images(data, '.pdf', result)
                return self.cleaner.clean_document(result)
        except Exception as e:
            print(f"⚠️ PDF fitz 解析失败, 回退 MarkItDown: {e}")
        fb = self._markitdown_convert(data, ".pdf")
        if not fb:
            return ""
        fb = self._upload_embedded_images(data, '.pdf', fb)
        return self.cleaner.clean_document(fb)

    # ── DOCX ─────────────────────────────────────────────────
    def _parse_docx(self, data: bytes, source_url: str = "") -> str:
        """DOCX: MarkItDown 优先 -> python-docx 回退(含图片) -> 清洗"""
        md_text = self._markitdown_convert(data, ".docx")
        if md_text:
            md_text = self._upload_embedded_images(data, '.docx', md_text)
            return self.cleaner.clean_document(md_text)

        if DocxDocument is None:
            return ""
        try:
            doc = DocxDocument(BytesIO(data))

            images = EmbeddedImageUploader.extract_from_zip(data, 'word/media/', min_size=self.MIN_IMG_BYTES)
            url_map: Dict[str, str] = {}
            if images:
                images.sort(key=lambda x: x[0])
                print(f"  📷 python-docx 回退: 从 DOCX 提取到 {len(images)} 张图片，正在上传...")
                url_map = EmbeddedImageUploader.upload_images(images)
                if url_map:
                    print(f"  ✅ 成功上传 {len(url_map)}/{len(images)} 张")

            img_count = 0
            rId_to_url: Dict[str, str] = {}
            for rel_id, rel in doc.part.rels.items():
                if "image" in getattr(rel, 'reltype', ''):
                    target = os.path.basename(str(rel.target_ref))
                    if target.lower() in {k.lower() for k in url_map}:
                        for k, v in url_map.items():
                            if k.lower() == target.lower():
                                rId_to_url[rel_id] = v
                                break

            paragraphs = []
            for para in doc.paragraphs:
                para_xml = para._element.xml
                has_image = '<w:drawing' in para_xml or '<v:imagedata' in para_xml or '<wp:inline' in para_xml
                text = para.text.strip()
                if text:
                    if para.style and para.style.name and 'Heading' in para.style.name:
                        level = para.style.name.replace('Heading', '').strip()
                        prefix = '#' * (int(level) if level.isdigit() else 2)
                        paragraphs.append(f"{prefix} {text}")
                    else:
                        paragraphs.append(text)
                if has_image:
                    img_count += 1
                    img_url = None
                    embed_match = re.search(r'r:embed="([^"]+)"', para_xml)
                    if embed_match:
                        img_url = rId_to_url.get(embed_match.group(1))
                    if not img_url and url_map:
                        ordered = sorted(url_map.values())
                        idx = img_count - 1
                        if idx < len(ordered):
                            img_url = ordered[idx]
                    if img_url:
                        paragraphs.append(f"![文档图片{img_count}]({img_url})")

            for table in doc.tables:
                rows = []
                for row in table.rows:
                    rows.append([cell.text.strip() for cell in row.cells])
                if rows:
                    paragraphs.append(self._rows_to_md_table(rows))

            md_text = "\n\n".join(paragraphs)
            print(f"  📄 MarkItDown 失败, python-docx 回退成功 ({len(md_text)} chars)")
        except Exception as e:
            print(f"⚠️ python-docx 回退也失败: {e}")
            return ""

        if not md_text:
            return ""
        return self.cleaner.clean_document(md_text)

    # ── PPTX ─────────────────────────────────────────────────
    def _parse_pptx(self, data: bytes, source_url: str = "") -> str:
        if PptxPresentation is not None:
            try:
                prs = PptxPresentation(BytesIO(data))
                parts = []
                img_count = 0
                images_to_upload: List[tuple] = []

                for si, slide in enumerate(prs.slides):
                    slide_title = ""
                    elements = []
                    for shape in slide.shapes:
                        top = shape.top or 0
                        if shape.has_text_frame:
                            text = "\n".join(
                                p.text.strip() for p in shape.text_frame.paragraphs if p.text.strip()
                            )
                            if text:
                                if not slide_title:
                                    try:
                                        if shape.is_placeholder and shape.placeholder_format.idx == 0:
                                            slide_title = text
                                    except Exception:
                                        pass
                                elements.append((top, text))
                        if PptxShapeType and shape.shape_type == PptxShapeType.PICTURE:
                            img_count += 1
                            placeholder = f"__PPTX_IMG_{img_count}__"
                            try:
                                blob = shape.image.blob
                                ext = getattr(shape.image, 'ext', 'png') or 'png'
                                if len(blob) >= self.MIN_IMG_BYTES:
                                    fname = f"pptx_s{si + 1}_img{img_count}.{ext}"
                                    mime = f"image/{ext}"
                                    images_to_upload.append((placeholder, fname, blob, mime))
                            except Exception:
                                pass
                            elements.append((top, f"![幻灯片{si + 1}-图片{img_count}]({placeholder})"))
                        if shape.has_table:
                            rows_data = []
                            for row in shape.table.rows:
                                rows_data.append([cell.text.strip() for cell in row.cells])
                            if rows_data:
                                elements.append((top, self._rows_to_md_table(rows_data)))
                    elements.sort(key=lambda x: x[0])
                    header = f"## 幻灯片 {si + 1}"
                    if slide_title:
                        header += f": {slide_title}"
                    body = "\n\n".join(e[1] for e in elements)
                    if body.strip():
                        parts.append(f"{header}\n\n{body}")

                if parts:
                    md_text = "\n\n---\n\n".join(parts)
                    if images_to_upload:
                        upload_list = [(fn, bl, mi) for _, fn, bl, mi in images_to_upload]
                        print(f"  📷 从 PPTX 幻灯片提取到 {len(upload_list)} 张图片，正在上传...")
                        url_map = EmbeddedImageUploader.upload_images(upload_list)
                        if url_map:
                            print(f"  ✅ 成功上传 {len(url_map)}/{len(upload_list)} 张")
                            for ph, fn, _, _ in images_to_upload:
                                if fn in url_map:
                                    md_text = md_text.replace(f"]({ph})", f"]({url_map[fn]})")
                    return self.cleaner.clean_document(md_text)
            except Exception as e:
                print(f"⚠️ python-pptx 解析失败, 回退 MarkItDown: {e}")
        fb = self._markitdown_convert(data, ".pptx")
        if not fb:
            return ""
        fb = self._upload_embedded_images(data, '.pptx', fb)
        return self.cleaner.clean_document(fb)

    # ── Excel (xlsx / xls) ───────────────────────────────────
    def _parse_excel(self, data: bytes, source_url: str = "") -> str:
        is_xls = not data[:4] == b'PK\x03\x04'
        suffix = ".xls" if is_xls else ".xlsx"
        md_text = self._markitdown_convert(data, suffix)
        if md_text and md_text.strip():
            return self.cleaner.clean_table(md_text)
        if not is_xls and openpyxl_load_workbook is not None:
            try:
                wb = openpyxl_load_workbook(BytesIO(data), read_only=True, data_only=True)
                parts = []
                for name in wb.sheetnames:
                    ws = wb[name]
                    rows = []
                    for ri, row in enumerate(ws.iter_rows(values_only=True)):
                        if ri >= self.MAX_TABLE_ROWS:
                            rows.append(["...", f"共 {ws.max_row} 行，已截断", "..."])
                            break
                        rows.append([str(c) if c is not None else "" for c in row])
                    if rows:
                        parts.append(f"### 工作表: {name}\n\n{self._rows_to_md_table(rows)}")
                wb.close()
                if parts:
                    return self.cleaner.clean_table("\n\n".join(parts))
            except Exception as e:
                print(f"⚠️ openpyxl 失败: {e}")
        if is_xls and xlrd is not None:
            try:
                wb = xlrd.open_workbook(file_contents=data)
                parts = []
                for name in wb.sheet_names():
                    ws = wb.sheet_by_name(name)
                    rows = []
                    for ri in range(min(ws.nrows, self.MAX_TABLE_ROWS)):
                        rows.append([str(ws.cell_value(ri, ci)) for ci in range(ws.ncols)])
                    if ws.nrows > self.MAX_TABLE_ROWS:
                        rows.append(["...", f"共 {ws.nrows} 行，已截断", "..."])
                    if rows:
                        parts.append(f"### 工作表: {name}\n\n{self._rows_to_md_table(rows)}")
                if parts:
                    return self.cleaner.clean_table("\n\n".join(parts))
            except Exception as e:
                print(f"⚠️ xlrd 失败: {e}")
        return ""

    # ── CSV ──────────────────────────────────────────────────
    def _parse_csv(self, data: bytes, source_url: str = "") -> str:
        md_text = self._markitdown_convert(data, ".csv")
        if md_text and md_text.strip():
            return self.cleaner.clean_table(md_text)
        text = self._decode_bytes(data)
        if not text:
            return ""
        try:
            dialect = csv.Sniffer().sniff(text[:8192])
            reader = csv.reader(text.splitlines(), dialect)
        except csv.Error:
            reader = csv.reader(text.splitlines())
        rows = []
        for i, row in enumerate(reader):
            if i >= self.MAX_TABLE_ROWS:
                rows.append(["...", "[已截断]", "..."])
                break
            rows.append(row)
        return self.cleaner.clean_table(self._rows_to_md_table(rows)) if rows else ""

    # ── HTML (文件) ──────────────────────────────────────────
    def _parse_html_file(self, data: bytes, source_url: str = "") -> str:
        text = self._decode_bytes(data)
        if not text:
            return ""
        cfg = use_config()
        cfg.set("DEFAULT", "EXTRACTION_TIMEOUT", "10")
        result = trafilatura.extract(
            text, config=cfg, output_format='markdown',
            include_images=True, favor_recall=True
        )
        return self.cleaner.clean_html(result) if result else ""

    # ── JSON ─────────────────────────────────────────────────
    def _parse_json(self, data: bytes, source_url: str = "") -> str:
        text = self._decode_bytes(data)
        if not text:
            return ""
        try:
            obj = json.loads(text)
            formatted = json.dumps(obj, ensure_ascii=False, indent=2)
        except json.JSONDecodeError:
            formatted = text
        if len(formatted) > self.MAX_JSON_CHARS:
            formatted = formatted[:self.MAX_JSON_CHARS] + "\n... [JSON 内容过长，已截断]"
        return f"```json\n{formatted}\n```"

    # ── XML ──────────────────────────────────────────────────
    def _parse_xml(self, data: bytes, source_url: str = "") -> str:
        md_text = self._markitdown_convert(data, ".xml")
        if md_text and md_text.strip():
            return self.cleaner.clean_text(md_text)
        text = self._decode_bytes(data)
        if not text:
            return ""
        if len(text) > self.MAX_TEXT_CHARS:
            text = text[:self.MAX_TEXT_CHARS] + "\n... [XML 过长，已截断]"
        return f"```xml\n{text}\n```"

    # ── Plain Text ───────────────────────────────────────────
    def _parse_plain_text(self, data: bytes, source_url: str = "") -> str:
        text = self._decode_bytes(data)
        return self.cleaner.clean_text(text) if text else ""

    # ── Markdown ─────────────────────────────────────────────
    def _parse_markdown(self, data: bytes, source_url: str = "") -> str:
        text = self._decode_bytes(data)
        return self.cleaner.clean_text(text) if text else ""

    # ── Image ────────────────────────────────────────────────
    def _parse_image(self, data: bytes, source_url: str = "") -> str:
        parts = []
        if source_url:
            parts.append(f"![image]({source_url})")
        if PILImage is not None:
            try:
                img = PILImage.open(BytesIO(data))
                w, h = img.size
                fmt = img.format or "Unknown"
                parts.append(f"**图片信息**: {fmt}, {w}x{h}px, {img.mode}")
                try:
                    import pytesseract
                    ocr_text = pytesseract.image_to_string(img, lang='chi_sim+eng')
                    if ocr_text and ocr_text.strip():
                        parts.append(f"\n**OCR 识别文本**:\n\n{ocr_text.strip()}")
                except (ImportError, Exception):
                    pass
            except Exception as e:
                print(f"⚠️ 图片解析失败: {e}")
        if not parts:
            return f"[图片文件, {len(data)} bytes]"
        return "\n\n".join(parts)

    # ── 工具方法 ─────────────────────────────────────────────
    @staticmethod
    def _decode_bytes(data: bytes) -> str:
        for enc in ('utf-8', 'utf-8-sig', 'gbk', 'gb2312', 'gb18030', 'big5', 'latin-1'):
            try:
                return data.decode(enc)
            except (UnicodeDecodeError, LookupError):
                continue
        return data.decode('utf-8', errors='replace')

    @staticmethod
    def _rows_to_md_table(rows: list) -> str:
        if not rows:
            return ""
        max_cols = max(len(r) for r in rows)
        padded = [r + [""] * (max_cols - len(r)) for r in rows]
        header = "| " + " | ".join(str(c).replace("|", "\\|").replace("\n", " ")[:80] for c in padded[0]) + " |"
        sep = "| " + " | ".join("---" for _ in padded[0]) + " |"
        body = []
        for row in padded[1:]:
            body.append("| " + " | ".join(str(c).replace("|", "\\|").replace("\n", " ")[:80] for c in row) + " |")
        return "\n".join([header, sep] + body)

    # ── 嵌入图片: 提取 + 上传 + 替换 ─────────────────────────
    _IMG_LOCAL_REF = re.compile(r'(!\[[^\]]*\])\((?!https?://|data:)([^)]+)\)')

    def _strip_and_replace_data_uris(self, md_text: str, ordered_urls: List[tuple]) -> str:
        """
        扫描 markdown, 把所有 ![alt](data:image/...;base64,...) 替换为上传后的真实 URL。
        完全不解码 base64 — 只用字符串定位 data: 开头和 ) 结尾, 然后整段替换。
        ordered_urls: [(filename, url), ...] 按文档中图片出现顺序排列。
        """
        url_idx = 0
        parts: List[str] = []
        pos = 0
        while pos < len(md_text):
            img_start = md_text.find('![', pos)
            if img_start == -1:
                parts.append(md_text[pos:])
                break
            bracket_close = md_text.find('](', img_start + 2)
            if bracket_close == -1:
                parts.append(md_text[pos:])
                break
            uri_start = bracket_close + 2
            if not md_text[uri_start:uri_start + 5] == 'data:':
                parts.append(md_text[pos:uri_start])
                pos = uri_start
                continue
            paren_close = md_text.find(')', uri_start)
            if paren_close == -1:
                parts.append(md_text[pos:])
                break
            alt = md_text[img_start + 2:bracket_close]
            parts.append(md_text[pos:img_start])
            if url_idx < len(ordered_urls):
                fname, url = ordered_urls[url_idx]
                parts.append(f"![{alt or fname}]({url})")
                url_idx += 1
            else:
                parts.append(f"![{alt or '图片'}]")
            pos = paren_close + 1
        return ''.join(parts)

    def _upload_embedded_images(self, data: bytes, ext: str, md_text: str) -> str:
        """
        核心逻辑: 从文档二进制直接提取图片 -> 上传到服务器获取真实 URL ->
        替换 markdown 中所有 data:image base64 引用和本地文件名引用。
        绝不解码 base64, 图片来源是文档 ZIP/PDF 二进制本身。
        """
        if not md_text:
            return md_text

        # ── Step 1: 从文档二进制提取真实图片文件 ──────────────────
        images: List[tuple] = []
        if ext == '.pdf':
            images = EmbeddedImageUploader.extract_from_pdf(
                data, max_pages=self.PDF_MAX_PAGES,
                min_size=self.MIN_IMG_BYTES, min_dim=self.MIN_IMG_DIM
            )
        elif ext in ('.docx', '.doc'):
            images = EmbeddedImageUploader.extract_from_zip(data, 'word/media/', min_size=self.MIN_IMG_BYTES)
        elif ext in ('.pptx', '.ppt'):
            images = EmbeddedImageUploader.extract_from_zip(data, 'ppt/media/', min_size=self.MIN_IMG_BYTES)

        if not images:
            return self._strip_and_replace_data_uris(md_text, [])

        images.sort(key=lambda x: x[0])

        # ── Step 2: 批量上传到服务器 ─────────────────────────────
        print(f"  📷 从文档提取到 {len(images)} 张图片，正在上传...")
        url_map = EmbeddedImageUploader.upload_images(images)
        if not url_map:
            print(f"  ⚠️ 图片上传失败，移除 base64 噪音")
            return self._strip_and_replace_data_uris(md_text, [])

        ordered_urls = [(fname, url_map[fname]) for fname, _, _ in images if fname in url_map]
        print(f"  ✅ 成功上传 {len(ordered_urls)}/{len(images)} 张图片")

        # ── Step 3: 替换 data:image base64 引用 (按顺序匹配) ────
        md_text = self._strip_and_replace_data_uris(md_text, ordered_urls)

        # ── Step 4: 替换本地文件名引用 ![](image1.png) (按名称匹配)
        local_matches = list(self._IMG_LOCAL_REF.finditer(md_text))
        if not local_matches:
            return md_text

        url_map_lower = {k.lower(): v for k, v in url_map.items()}
        name_no_ext_map = {os.path.splitext(k)[0].lower(): v for k, v in url_map.items()}

        for m_obj in reversed(local_matches):
            prefix = m_obj.group(1)
            ref = m_obj.group(2)
            ref_base = os.path.basename(ref).lower()
            ref_no_ext = os.path.splitext(ref_base)[0]
            new_url = (
                    url_map_lower.get(ref_base)
                    or url_map_lower.get(ref.lower())
                    or name_no_ext_map.get(ref_no_ext)
                    or name_no_ext_map.get(ref.lower())
            )
            if new_url:
                replacement = f"{prefix}({new_url})"
                md_text = md_text[:m_obj.start()] + replacement + md_text[m_obj.end():]

        return md_text

    # ── 主入口 (同步, 在 asyncio.to_thread 中调用) ───────────
    def parse(self, binary_content: bytes, file_extension: str, source_url: str = "") -> str:
        ext = file_extension.lower().strip()
        if not ext.startswith("."):
            ext = f".{ext}"

        result = ""
        if ext == '.pdf':
            result = self._parse_pdf(binary_content, source_url)
        elif ext == '.docx':
            result = self._parse_docx(binary_content, source_url)
        elif ext == '.pptx':
            result = self._parse_pptx(binary_content, source_url)
        elif ext in ('.xlsx', '.xls'):
            result = self._parse_excel(binary_content, source_url)
        elif ext == '.csv':
            result = self._parse_csv(binary_content, source_url)
        elif ext in ('.html', '.htm'):
            result = self._parse_html_file(binary_content, source_url)
        elif ext == '.json':
            result = self._parse_json(binary_content, source_url)
        elif ext == '.xml':
            result = self._parse_xml(binary_content, source_url)
        elif ext == '.txt':
            result = self._parse_plain_text(binary_content, source_url)
        elif ext in ('.md', '.markdown'):
            result = self._parse_markdown(binary_content, source_url)
        elif ext in ('.jpg', '.jpeg', '.png', '.gif', '.webp', '.bmp'):
            result = self._parse_image(binary_content, source_url)
        elif ext in ('.doc', '.ppt'):
            fb = self._markitdown_convert(binary_content, ext)
            if fb:
                fb = self._upload_embedded_images(binary_content, ext, fb)
                result = self.cleaner.clean_document(fb)

        if not (result and result.strip()):
            print(f"⚠️ 格式 {ext} 专用解析无结果，尝试 MarkItDown 兜底")
            fb = self._markitdown_convert(binary_content, ext)
            if fb:
                if ext in ('.pdf', '.docx', '.doc', '.pptx', '.ppt'):
                    fb = self._upload_embedded_images(binary_content, ext, fb)
                result = self.cleaner.clean_document(fb)
            else:
                result = f"[无法解析 {ext} 格式文件]"

        return result

    async def parse_async(self, binary_content: bytes, file_extension: str, source_url: str = "") -> str:
        return await asyncio.to_thread(self.parse, binary_content, file_extension, source_url)

    def parse_html_content(self, html: str, base_url: str = "") -> str:
        if not html:
            return ""
        cfg = use_config()
        cfg.set("DEFAULT", "EXTRACTION_TIMEOUT", "10")
        result = trafilatura.extract(
            html, config=cfg, output_format='markdown',
            include_images=True, favor_recall=True
        )
        return self.cleaner.clean_html(result) if result else ""


class ContentScraper(ABC):
    """内容抓取器的抽象基类。"""

    # 【修改】方法签名，接收一个包含所有元数据的字典
    @abstractmethod
    async def scrape(self, item_info: Dict[str, Any], client: httpx.AsyncClient) -> Dict[str, Any]:
        """
        抓取单个URL的内容，并返回合并了原始信息的结果。
        """
        pass


# --- 2.1 SearchAPI.io 的手动抓取与清洗实现 ---
class SearchApiScraper(ContentScraper):
    def __init__(self):
        self.trafilatura_config = use_config()
        self.trafilatura_config.set("DEFAULT", "EXTRACTION_TIMEOUT", "10")
        self.parser_service = DocumentParserService()

        # 编译常用的正则表达式以提高性能
        self.NOISY_PATTERNS = [re.compile(p, re.IGNORECASE) for p in [
            r'^\s*$', r'^[\-=*#_]{3,}$', r'.*\.(html|shtml|htm|php)\s*$',
            r'.{0,50}(搜狐|网易|腾讯|新浪|登录|注册|版权所有|版权声明).{0,50}$',
            r'\[\d+\]|\[下一页\]|\[上一页\]', r'\[(编辑|查看历史|讨论|阅读|来源|原标题)\]',
            r'^\*+\s*\[.*?\]\(.*?\)',
            r'^\s*(分享到|扫描二维码|返回搜狐|查看更多|责任编辑|记者|通讯员)',
            r'^\s*([京公网安备京网文京ICP备]|互联网新闻信息服务许可证|信息网络传播视听节目许可证)',
        ]]
        self.IMG_PATTERN = re.compile(r'(!\[(.*?)\]\((.*?)\))')
        self.LINK_PATTERN = re.compile(r'\[.*?\]\(.*?\)')
        self.EDITOR_PATTERN = re.compile(r'(\(|\[)\s*责任编辑：.*?\s*(\)|\])')

    # --- 2.1.1 内容提取工具 (来自您的代码) ---
    def _extract_pdf_text(self, binary_content: bytes) -> str:
        """
        使用 PyMuPDF (fitz) 从 PDF 的二进制内容中快速提取文本。
        """
        text_parts = []
        try:
            # 直接从内存中的字节流打开 PDF
            with fitz.open(stream=binary_content, filetype="pdf") as doc:
                # 限制处理的页数，避免处理超大文件
                num_pages_to_process = min(len(doc), self.PDF_MAX_PAGES_TO_PROCESS)
                if len(doc) > self.PDF_MAX_PAGES_TO_PROCESS:
                    print(f"  📄 PDF 页数过多 ({len(doc)} pages), 只处理前 {self.PDF_MAX_PAGES_TO_PROCESS} 页。")
                for i in range(num_pages_to_process):
                    page = doc.load_page(i)
                    page_text = page.get_text("text", sort=True)  # sort=True 尝试保持阅读顺序
                    if page_text:
                        text_parts.append(page_text)

            return "\n\n".join(text_parts).strip()
        except Exception as e:
            print(f"⚠️ PyMuPDF (fitz) 解析失败: {e}")
            return ""

    def _parse_videos_from_html(self, html: str, base_url: str) -> List[str]:
        try:
            soup = BeautifulSoup(html, "lxml")
            videos = []
            for video in soup.find_all("video"):
                src = video.get("src")
                if src: videos.append(urljoin(base_url, src))
                for source in video.find_all("source"):
                    src = source.get("src")
                    if src: videos.append(urljoin(base_url, src))
            for iframe in soup.find_all("iframe"):
                src = iframe.get("src")
                if src and any(k in src for k in ["youtube", "vimeo", "embed", ".mp4"]):
                    videos.append(urljoin(base_url, src))
            return list(dict.fromkeys(videos))  # 去重并保持顺序
        except Exception as e:
            print(f"⚠️ 视频解析失败: {e}")
            return []

    # --- 2.1.2 内容清洗工具 (来自您的代码，已优化和异步化) ---
    async def _is_valid_image_url_async(self, url: str, client: httpx.AsyncClient) -> bool:
        if not url or not url.startswith(('http://', 'https://')): return False
        try:
            resp = await client.head(url, timeout=5, follow_redirects=True)
            content_type = resp.headers.get('content-type', '').lower()
            return resp.is_success and 'image' in content_type
        except httpx.RequestError:
            return False

    async def _remove_invalid_images_async(self, md: str, client: httpx.AsyncClient) -> str:
        MAX_IMAGES_TO_VALIDATE = 25
        matches = list(self.IMG_PATTERN.finditer(md))
        urls_to_check_all = {m.group(3).strip() for m in matches}

        urls_to_check = set(list(urls_to_check_all)[:MAX_IMAGES_TO_VALIDATE])
        if len(urls_to_check_all) > MAX_IMAGES_TO_VALIDATE:
            print(f"⚠️ 图片数量过多 ({len(urls_to_check_all)}), 只验证前 {MAX_IMAGES_TO_VALIDATE} 张。")

        tasks = {url: self._is_valid_image_url_async(url, client) for url in urls_to_check}
        results = await asyncio.gather(*tasks.values(), return_exceptions=True)

        url_status = dict(zip(tasks.keys(), results))
        valid_urls = {url for url, res in url_status.items() if isinstance(res, bool) and res}
        valid_urls.update(urls_to_check_all - urls_to_check)  # 未检查的默认有效

        def replacer(match: re.Match):
            return match.group(0) if match.group(3).strip() in valid_urls else ""

        return self.IMG_PATTERN.sub(replacer, md)

    def _is_noisy_line(self, line: str) -> bool:
        stripped = line.strip()
        for pat in self.NOISY_PATTERNS:
            if pat.search(stripped): return True
        links = self.LINK_PATTERN.findall(stripped)
        if len(links) > 2 and len(stripped) / (len(links) + 1) < 30: return True
        return False

    async def _clean_content_async(self, text: str, client: httpx.AsyncClient) -> str:
        if not text: return ""
        text = await self._remove_invalid_images_async(text, client)

        lines = text.splitlines()
        cleaned_lines = []
        for line in lines:
            if not self._is_noisy_line(line):
                line = self.EDITOR_PATTERN.sub('', line).strip()
                if line: cleaned_lines.append(line)

        # 去除连续空行
        out = []
        for i, line in enumerate(cleaned_lines):
            if i > 0 and not line.strip() and not cleaned_lines[i - 1].strip():
                continue
            out.append(line)

        return "\n".join(out).strip()

    # --- 2.1.3 主抓取函数 (来自您的代码，封装为scrape方法) ---
    async def scrape(self, item_info: Dict[str, Any], client: httpx.AsyncClient) -> dict:
        url = item_info.get("url")
        print(f"🕸️ [SearchAPI Scraper] 开始处理: {url}")
        headers = {
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36'}

        try:
            # 1. 检测是否为支持的文档类型
            supported_extensions = {
                ".pdf", ".docx", ".doc", ".pptx", ".ppt", ".xlsx", ".xls", ".csv",
                ".txt", ".md", ".markdown", ".json", ".xml",
                ".jpg", ".jpeg", ".png", ".gif", ".webp", ".bmp",
            }
            url_lower = url.lower()
            ext = os.path.splitext(url_lower)[1]

            is_document = ext in supported_extensions
            content_type = ""

            # 如果这类已知扩展名，可以直接认定为文档，减少 HEAD 请求（或保留 HEAD 以检测大小）
            # 这里保留 HEAD 请求以获取 Content-Type 和 Content-Length
            async with client.stream("HEAD", url, headers=headers, follow_redirects=True) as head_response:
                content_type = head_response.headers.get('content-type', '').lower()
                content_length = int(head_response.headers.get('content-length', 0))

                # Check for document mime types if extension wasn't obvious
                if not is_document:
                    if 'pdf' in content_type:
                        is_document = True;
                        ext = ".pdf"
                    elif 'word' in content_type or 'officedocument' in content_type:
                        is_document = True;
                        ext = ".docx"  # Simplification
                    elif 'excel' in content_type or 'spreadsheet' in content_type:
                        is_document = True;
                        ext = ".xlsx"
                    elif 'powerpoint' in content_type or 'presentation' in content_type:
                        is_document = True;
                        ext = ".pptx"
                    elif 'csv' in content_type:
                        is_document = True;
                        ext = ".csv"
                    elif 'image/' in content_type:
                        is_document = True
                        if 'png' in content_type:
                            ext = ".png"
                        elif 'gif' in content_type:
                            ext = ".gif"
                        elif 'webp' in content_type:
                            ext = ".webp"
                        elif 'bmp' in content_type:
                            ext = ".bmp"
                        else:
                            ext = ".jpg"
                    elif 'application/json' in content_type:
                        is_document = True;
                        ext = ".json"
                    elif 'text/xml' in content_type or 'application/xml' in content_type:
                        is_document = True;
                        ext = ".xml"
                    elif 'text/plain' in content_type and ext not in supported_extensions:
                        is_document = True;
                        ext = ".txt"
                    elif 'text/markdown' in content_type:
                        is_document = True;
                        ext = ".md"

                # 【优化】检查文件大小 (限制为 20MB)
                MAX_SIZE = 20 * 1024 * 1024
                if is_document and content_length > MAX_SIZE:
                    raise ValueError(f"文档文件过大 ({content_length / 1024 / 1024:.2f}MB > 20MB)，跳过处理。")

            # 根据类型执行不同逻辑
            raw_content, final_url = "", url

            if is_document:
                print(f"  📄 [SearchAPI Scraper] 检测到文档 ({ext}): {url}")
                try:
                    async def _download_doc():
                        resp = await client.get(url, timeout=None, headers=headers, follow_redirects=True)
                        resp.raise_for_status()
                        return str(resp.url), await resp.aread()

                    final_url, file_bytes = await asyncio.wait_for(_download_doc(), timeout=60)
                except asyncio.TimeoutError:
                    raise TimeoutError(f"下载超时 (60s): {url}")

                print(f"  🚀 [SearchAPI Scraper] 正在使用 DocumentParserService 解析 ({ext})...")
                raw_content = await asyncio.to_thread(
                    self.parser_service.parse, file_bytes, ext or ".bin", url
                )
            else:
                print(f"  📑 [SearchAPI Scraper] 检测到 HTML: {url}")
                response = await client.get(url, timeout=20, headers=headers, follow_redirects=True)
                response.raise_for_status()
                final_url = str(response.url)
                html_content = response.text
                raw_content = await asyncio.to_thread(
                    trafilatura.extract, html_content, config=self.trafilatura_config,
                    output_format='markdown', include_images=True, favor_recall=True)

                # HTML的视频解析和清洗
                if raw_content:
                    print(f"  🧹 [SearchAPI Scraper] 正在清洗HTML内容: {final_url}")
                    cleaned_content = await self._clean_content_async(raw_content, client)
                    print(cleaned_content)
                    videos = self._parse_videos_from_html(html_content, final_url)
                    if videos:
                        video_section = "\n\n## 参考视频:\n" + "\n".join(f"- {vid}" for vid in videos)
                        cleaned_content += video_section
                    raw_content = cleaned_content

            if not raw_content: raise ValueError("内容提取返回为空。")
            print(f"✅ [SearchAPI Scraper] 成功: {url}")
            return {**item_info, "url": final_url, "content": raw_content, "status": "success"}
        except Exception as e:
            error_msg = f"处理失败 {url}: {type(e).__name__} - {e}"
            print(f"⚠️ [SearchAPI Scraper] {error_msg}")
            return {**item_info, "content": "", "status": "failed", "error_message": str(e)}


# --- 2.2 FirecrawlScraper ---
class FirecrawlScraper(ContentScraper):
    def __init__(self):
        self.api_key = os.environ.get("FIRECRAWL_API_KEY", "fc-a36b7d2fb273485680d0fe6abd686935")
        if not self.api_key: raise ValueError("未提供 Firecrawl API Key。")
        self.base_url = "https://api.firecrawl.dev/v2/scrape"
        self.headers = {"Authorization": f"Bearer {self.api_key}", "Content-Type": "application/json"}

    async def scrape(self, item_info: Dict[str, Any], client: httpx.AsyncClient) -> dict:
        url = item_info.get("url")
        print(f"🔥 [Firecrawl Scraper] 开始处理: {url}")
        try:
            # 【修改】根据文档，移除 pageOptions，将选项置于顶层
            payload = {
                "url": url,
                "onlyMainContent": True,
                "removeBase64Images": True,
                "blockAds": True
            }
            resp = await client.post(self.base_url, headers=self.headers, json=payload, timeout=45)

            if not resp.is_success:
                try:
                    error_details = resp.json()
                    raise httpx.HTTPStatusError(f"API返回错误: {error_details.get('error', str(error_details))}",
                                                request=resp.request, response=resp)
                except json.JSONDecodeError:
                    resp.raise_for_status()

            data_wrapper = resp.json()

            # 【修改】根据文档，检查顶层 success 键和 data 字段
            if not data_wrapper.get("success"):
                raise ValueError(f"API返回失败状态: {data_wrapper.get('error', '未知错误')}")

            data = data_wrapper.get("data")
            if not data:
                raise ValueError("API返回的 'data' 字段为空。")

            content = data.get("markdown")
            if content is None:
                raise ValueError("API未返回 'markdown' 字段。")

            final_url = data.get("metadata", {}).get("sourceURL", url)

            print(f"✅ [Firecrawl Scraper] 成功: {url}")
            return {**item_info, "url": final_url, "content": content, "status": "success"}

        except Exception as e:
            error_msg = f"处理失败 {url}: {type(e).__name__} - {e}"
            print(f"⚠️ [Firecrawl Scraper] {error_msg}")
            return {**item_info, "content": "", "status": "failed", "error_message": str(e)}


# --- 2.3 JinaScraper ---
class JinaScraper(ContentScraper):
    def __init__(self):
        self.api_key = os.environ.get("JINA_API_KEY",
                                      "jina_b4348ffc39ca47bfbe753b95f59428c7i6ifkOFXRPdF3dRa5Rwb6T8FvrLH")
        if not self.api_key: raise ValueError("未提供 Jina API Key。")
        self.base_url = "https://r.jina.ai/"
        self.headers = {
            "Authorization": f"Bearer {self.api_key}",
            "Content-Type": "application/json",
            "Accept": "application/json",
            "X-Return-Format": "markdown"  # 关键：直接获取Markdown
        }

    async def scrape(self, item_info: Dict[str, Any], client: httpx.AsyncClient) -> dict:
        url = item_info.get("url")
        print(f"🌀 [Jina Scraper] 开始处理: {url}")
        try:
            # Jina的Reader API 对GET请求更友好，直接拼接URL
            target_url = f"{self.base_url}{url}"
            resp = await client.get(target_url, headers=self.headers, timeout=45)
            resp.raise_for_status()

            # 【修改】根据文档，Jina 可能直接返回 Markdown 文本，也可能返回 JSON
            content_type = resp.headers.get("content-type", "").lower()
            if "application/json" in content_type:
                data_wrapper = resp.json()
                if data_wrapper.get("code") == 200 and "data" in data_wrapper:
                    data = data_wrapper["data"]
                    content = data.get("content")
                    final_url = data.get("url", url)
                    if content is None: raise ValueError("API JSON响应中缺少 'content' 字段。")
                else:
                    raise ValueError(f"API JSON响应错误: {data_wrapper}")
            else:
                # 假设直接返回Markdown文本
                content = resp.text
                final_url = url

            if not content.strip(): raise ValueError("API 返回内容为空。")

            print(f"✅ [Jina Scraper] 成功: {url}")
            return {**item_info, "url": final_url, "content": content, "status": "success"}

        except Exception as e:
            error_msg = f"处理失败 {url}: {type(e).__name__} - {e}"
            print(f"⚠️ [Jina Scraper] {error_msg}")
            return {**item_info, "content": "", "status": "failed", "error_message": str(e)}


# --- 2.4 TavilyScraper ---
class TavilyScraper(ContentScraper):
    def __init__(self):
        self.api_key = os.environ.get("TAVILY_API_KEY", "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")
        if not self.api_key: raise ValueError("未提供 Tavily API Key。")
        self.base_url = "https://api.tavily.com/extract"
        self.headers = {"Authorization": f"Bearer {self.api_key}", "Content-Type": "application/json"}

    async def scrape(self, item_info: Dict[str, Any], client: httpx.AsyncClient) -> dict:
        url = item_info.get("url")
        print(f"🤖 [Tavily Scraper] 开始处理: {url}")
        try:
            # 【修改】根据文档，urls字段应该是列表，且使用 format: markdown
            payload = {"urls": [url], "format": "markdown"}
            resp = await client.post(self.base_url, json=payload, headers=self.headers, timeout=45)
            resp.raise_for_status()
            data = resp.json()

            if not data.get("results") or not isinstance(data["results"], list):
                failed_info = data.get("failed_results", [])
                raise ValueError(f"API调用失败: {failed_info}")

            result = data["results"][0]
            # 【修改】根据文档，内容字段是 raw_content
            content = result.get("raw_content")
            if content is None: raise ValueError("API未返回raw_content内容。")

            final_url = result.get("url", url)

            print(f"✅ [Tavily Scraper] 成功: {url}")
            return {**item_info, "url": final_url, "content": content, "status": "success"}

        except Exception as e:
            error_msg = f"处理失败 {url}: {type(e).__name__} - {e}"
            print(f"⚠️ [Tavily Scraper] {error_msg}")
            return {**item_info, "content": "", "status": "failed", "error_message": str(e)}


# --- 2.5 ZhiLianJobScraper ---
class ZhiLianJobScraper:
    def __init__(self):
        self.api_url = "http://119.45.167.133:12906/api/scrape/zhilian"
        self.headers = {'accept': 'application/json', 'Content-Type': 'application/json'}

    async def scrape_jobs(self, payload: Dict[str, Any], client: httpx.AsyncClient) -> Dict[str, Any]:
        print(f"💼 [ZhiLian Scraper] 开始使用负载调用API: {json.dumps(payload, ensure_ascii=False)}")
        if not payload or not payload.get("keywords") or not payload.get("provinces"):
            msg = "负载无效，缺少 'keywords' 或 'provinces'。"
            print(f"⚠️ [ZhiLian Scraper] {msg}")
            return {"status": "skipped", "data": [], "message": msg}
        try:
            # 确保 page_size 是整数
            if 'page_size' in payload: payload['page_size'] = int(payload['page_size'])

            resp = await client.post(self.api_url, headers=self.headers, json=payload, timeout=60)
            resp.raise_for_status()
            response_data = resp.json()
            if response_data.get("code") == 200:
                print(f"✅ [ZhiLian Scraper] 成功: {response_data.get('message')}")
                return {"status": "success", "data": response_data.get("data", []),
                        "message": response_data.get("message")}
            else:
                msg = f"API返回错误码 {response_data.get('code')}: {response_data.get('message')}"
                print(f"API returned non-200 code: {msg}")
                return {"status": "failed", "data": [], "message": msg}
        except Exception as e:
            error_msg = f"API请求失败: {type(e).__name__} - {e}"
            print(f"⚠️ [ZhiLian Scraper] {error_msg}")
            return {"status": "failed", "data": [], "message": error_msg}


class TianyanEnterpriseScraper:
    def __init__(self):
        self.api_url = "http://open.api.tianyancha.com/services/open/ic/baseinfo/normal"
        # 从环境变量或直接硬编码获取Token
        self.token = os.environ.get("TIANYANCHA_TOKEN", "4d882100-ed23-4c22-a83b-c77af2e4be42")
        self.headers = {'Authorization': self.token}

    async def scrape_enterprise(self, name: str, client: httpx.AsyncClient) -> Dict[str, Any]:
        print(f"🏢 [Tianyan Scraper] 开始查询企业: {name}")
        base_return = {"query_name": name}
        if not name:
            msg = "企业名称为空，跳过查询。"
            print(f"🟡 [Tianyan Scraper] {msg}")
            return {**base_return, "status": "skipped", "data": None, "message": msg}
        try:
            params = {"keyword": name}
            resp = await client.get(self.api_url, headers=self.headers, params=params, timeout=30)
            resp.raise_for_status()
            response_data = resp.json()
            if response_data.get("error_code") == 0:
                print(f"✅ [Tianyan Scraper] 成功查询到: {name}")
                return {**base_return, "status": "success", "data": response_data.get("result"),
                        "message": response_data.get("reason")}
            else:
                msg = f"API返回错误码 {response_data.get('error_code')}: {response_data.get('reason')}"
                print(f"⚠️ [Tianyan Scraper] {msg}")
                return {**base_return, "status": "failed", "data": None, "message": msg}
        except Exception as e:
            error_msg = f"API请求失败: {type(e).__name__} - {e}"
            print(f"⚠️ [Tianyan Scraper] {error_msg}")
            return {**base_return, "status": "failed", "data": None, "message": error_msg}


class DataOrchestrator:
    def __init__(self):
        self.content_scrapers: Dict[str, ContentScraper] = {
            "searchapi": SearchApiScraper(), "firecrawl": FirecrawlScraper(),
            "jina": JinaScraper(), "tavily": TavilyScraper(),
        }
        self.job_scraper = ZhiLianJobScraper()
        self.enterprise_scraper = TianyanEnterpriseScraper()

    # 【调整】整个 process_all 方法被重构，以实现条件化任务调度。
    async def process_all(
            self,
            web_url_info_list: List[Dict[str, Any]],
            career_payload: Dict,
            enterprise_names: List[str]  # 接收列表
    ) -> Dict[str, Any]:
        """
        根据有效的输入，条件化地创建并并发执行所有抓取任务。
        """
        final_results = {
            "content_results": [],
            "job_result": None,
            "enterprise_results": [],  # 默认返回空列表
        }
        ssl_context = httpx.create_ssl_context(verify=False)
        async with httpx.AsyncClient(http2=True, verify=ssl_context, timeout=30, follow_redirects=True,
                                     limits=httpx.Limits(max_connections=50)) as client:

            content_tasks = []
            if web_url_info_list:
                print(f"  [Orchestrator] 准备 {len(web_url_info_list)}个网页抓取任务。")
                for item in web_url_info_list:
                    scraper = self.content_scrapers.get(item.get("provider")) or self.content_scrapers["searchapi"]
                    content_tasks.append(scraper.scrape(item, client))
            job_tasks = []
            if career_payload and career_payload.get("keywords") and career_payload.get("provinces"):
                print("  [Orchestrator] 准备招聘信息抓取任务。")
                job_tasks.append(self.job_scraper.scrape_jobs(career_payload, client))
            else:
                print("  [Orchestrator] 招聘信息负载无效，跳过任务。")

            # 【调整】为列表中的每个企业名称创建查询任务
            enterprise_tasks = []
            if enterprise_names:
                print(f"  [Orchestrator] 准备 {len(enterprise_names)}个企业信息查询任务。")
                for name in enterprise_names:
                    enterprise_tasks.append(self.enterprise_scraper.scrape_enterprise(name, client))
            else:
                print("  [Orchestrator] 企业名称列表为空，跳过任务。")

            tasks_to_run = content_tasks + job_tasks + enterprise_tasks
            if not tasks_to_run:
                print("  [Orchestrator] 没有可执行的任务。")
                return final_results
            all_results = await asyncio.gather(*tasks_to_run, return_exceptions=True)
            # 【调整】安全地解析和分离三组任务的结果
            content_end_idx = len(content_tasks)
            job_end_idx = content_end_idx + len(job_tasks)
            final_results["content_results"] = all_results[:content_end_idx]

            job_task_results = all_results[content_end_idx:job_end_idx]
            if job_task_results:
                final_results["job_result"] = job_task_results[0]

            final_results["enterprise_results"] = all_results[job_end_idx:]
            return final_results


# --- 5. Dify 节点主入口 ---
async def main_async(raw_input: Any) -> Dict[str, Any]:
    # 1. 解析输入
    parsed_data = _parse_input_data(raw_input)
    web_url_info_list = parsed_data["web_url_info_list"]
    video_url_info_list = parsed_data["video_url_info_list"]
    career_payload = parsed_data["career_payload"]
    enterprise_names = parsed_data["enterprise_names"]
    mode = parsed_data.get("mode", "legacy")

    if not web_url_info_list and not video_url_info_list and (
            not career_payload or not career_payload.get("keywords")) and not enterprise_names:
        print("🟡 所有输入均为空，提前返回。")
        return {"scraped_datas": {}, "scraped_datas_str": "{}"}
    # 2. 运行调度器
    orchestrator = DataOrchestrator()
    results = await orchestrator.process_all(web_url_info_list, career_payload, enterprise_names)

    # 3. 格式化网页内容输出 (按 origin_key 分组)
    content_results_by_origin = {
        "comprehensive_data": [],
        "general_web_data": [],
        "institution_source_data": []
    }

    for i, result in enumerate(results["content_results"]):
        if isinstance(result, Exception): continue
        if result.get("status") == "success":
            sanitized_url = re.sub(r'[^a-zA-Z0-9]', '-',
                                   result.get("url", "").replace("https://", "").replace("http://", ""))

            # Retrieve the origin_key from the input list corresponding to this result
            # Note: results["content_results"] corresponds exactly to web_url_info_list order
            origin_key = web_url_info_list[i].get("origin_key", "comprehensive_data")

            item_data = {
                "type": "web", "source_id": f"web-{sanitized_url[:100]}", "url": result.get("url"),
                "title": result.get("title"), "source": result.get("source"), "snippet": result.get("snippet"),
                "query": result.get("query"), "content": result.get("content", "")
            }

            if origin_key in content_results_by_origin:
                content_results_by_origin[origin_key].append(item_data)
            else:
                # Fallback
                content_results_by_origin["comprehensive_data"].append(item_data)

    # 4. 格式化视频内容输出 (暂不分组，视频通常只出现在 comprehensive 或 general 中，这里简单处理)
    # 如果需要严格分组，也需要在 _extract_urls 中对视频加 origin_key
    all_video_list = []
    for video_item in video_url_info_list:
        all_video_list.append({
            "type": "video", "url": video_item.get("url"), "title": video_item.get("title"),
            "source": video_item.get("source"), "snippet": video_item.get("snippet"),
            "video_id": video_item.get("video_id"), "embed_url": video_item.get("embed_url"),
            "thumbnail_url": video_item.get("thumbnail_url"), "query": video_item.get("query")
        })

    # 【调整】处理招聘和企业信息结果时，检查它们是否存在（是否为None）
    career_postings = results.get("job_result")
    if career_postings is None:
        career_postings = {}  # 如果未执行，则返回空对象
    elif isinstance(career_postings, Exception):
        career_postings = {"status": "failed", "data": [], "message": f"任务异常: {career_postings}"}

    enterprise_infos_raw = results.get("enterprise_results", [])
    enterprise_infos_output = {}
    if enterprise_names:
        successful_data = []
        failed_queries = []

        for item in enterprise_infos_raw:
            res = item
            if isinstance(item, Exception):
                res = {"status": "failed", "message": f"任务执行异常: {str(item)}", "query_name": "Unknown"}

            if res.get("status") == "success" and res.get("data"):
                successful_data.append(res["data"])
            elif res.get("status") in ["failed", "skipped"]:
                failed_queries.append({
                    "query_name": res.get("query_name", "N/A"),
                    "error_message": res.get("message", "未知错误")
                })
        final_status = "skipped"
        if successful_data or failed_queries:
            if not failed_queries:
                final_status = "success"
            elif not successful_data:
                final_status = "failed"
            else:
                final_status = "partial_success"

        enterprise_infos_output = {
            "status": final_status,
            "data": successful_data,
            "failed_queries": failed_queries,
            "summary": f"共查询 {len(enterprise_names)} 个企业，成功 {len(successful_data)} 个，失败 {len(failed_queries)} 个。"
        }

    # 6. 组装最终输出
    final_output = {}

    if mode == "legacy":
        # Legacy Mode Output Structure
        comprehensive_data_output = {
            "all_source_list": content_results_by_origin["comprehensive_data"],
            "all_video_list": all_video_list
        }
        final_output = {
            "scraped_datas": {
                "comprehensive_data": comprehensive_data_output,
                "career_postings": career_postings,
                "enterprise_infos": enterprise_infos_output
            }
        }
    else:
        # External Mode Output Structure
        # Note: External mode currently doesn't focus on video list separately in the top structure, 
        # but we can include it if needed. For now, we follow the plan to separate general and institution.
        final_output = {
            "scraped_datas": {
                "general_web_data": content_results_by_origin["general_web_data"],
                "institution_source_data": content_results_by_origin["institution_source_data"],
                "career_postings": career_postings,
                # Video list can be appended to general_web_data or kept separate if required by downstream.
                # For now, let's keep it simple.
            }
        }

    return {
        "scraped_datas": final_output["scraped_datas"],
        "scraped_datas_str": json.dumps(final_output, ensure_ascii=False, indent=2)
    }


def main(datas_input: Any) -> Dict[str, Any]:
    try:
        return _dify_debug_return(asyncio.run(main_async(raw_input=datas_input)))
    except Exception as e:
        print(f"‼️ 节点执行时发生顶层错误: {e}")
        error_payload = {
            "comprehensive_data": {
                "all_source_list": [
                    {"type": "web", "source_id": "NODE_EXECUTION_ERROR", "title": "节点执行失败", "url": "",
                     "content": f"An error occurred: {str(e)}\n\n{traceback.format_exc()}"}],
                "all_video_list": []
            },
            "career_postings": {},
            "enterprise_infos": {
                "status": "failed",
                "data": [],
                "failed_queries": [{"query_name": "Node Execution", "error_message": "节点顶层异常"}],
                "summary": "节点执行失败"
            }
        }
        return _dify_debug_return({
            "scraped_datas": error_payload,
            "scraped_datas_str": json.dumps({"scraped_datas": error_payload}, ensure_ascii=False, indent=2)
        })

# main({'career_data': {},
#       'general_web_data': [{'errors': [],
#                             'query': '保育员',
#                             'web_results': [
#                                 {
#                                     'searchapi_snippet': '蓝盈莹在短剧《马背摇篮》中饰演保育员文纫秋，虽无实际育儿经验，却通过与136名小演员的真实互动、沉浸式重走历史路线及细腻的表演设计，将战时“文妈妈”的柔韧与坚毅演绎 '
#                                                          '...',
#                                     'searchapi_source': '新浪新闻_手机新浪网',
#                                     'searchapi_title': "演员蓝盈莹无子女却演活'文妈妈'，真实儿童互动如何成就短 ...",
#                                     'searchapi_type': 'web',
#                                     'searchapi_url': 'https://news.sina.cn/bignews/insight/2026-01-24/detail-inhikete1442916.d.html?oid=%E9%AB%98%E4%BB%BFmiumiu%E5%A5%B3%E5%8C%85%E7%B2%89%E7%BA%A2%E8%89%B2%EF%BC%88%E5%BE%AE%E4%BF%A1198099199%EF%BC%89lvN7&vt=4'},
#                                 {
#                                     'searchapi_snippet': '保育员是幼儿园重要工种之一，是保育工作的具体实施者。虽然每个幼儿园的工种有所差异但是其基本职责与要求都是一样的，但是其目的是促进幼儿的全面发展。',
#                                     'searchapi_source': 'orginview.com',
#                                     'searchapi_title': '幼儿园的保育员需要什么证（需要什么条件好不好做） -',
#                                     'searchapi_type': 'web',
#                                     'searchapi_url': 'http://www.orginview.com/plugin.php?id=tom_tctoutiao&site=1&mod=info&aid=398'},
#                                 {
#                                     'searchapi_snippet': '1'
#                                                          '...',
#                                     'searchapi_source': 'ppt',
#                                     'searchapi_title': "演员蓝盈莹无子女却演活'文妈妈'，真实儿童互动如何成就短 ...",
#                                     'searchapi_type': 'web',
#                                     'searchapi_url': 'https://www.sem.tsinghua.edu.cn/__local/6/42/E9/AD5A763C99A09DD5A0C0C148278_E177A355_78F425.pptx?e=.pptx'},
#                                 {
#                                     'searchapi_snippet': '2',
#                                     'searchapi_source': 'word',
#                                     'searchapi_title': '幼儿园的保育员需要什么证（需要什么条件好不好做） -',
#                                     'searchapi_type': 'web',
#                                     'searchapi_url': 'https://www.hbea.edu.cn/files/2024-05/%E9%99%84%E4%BB%B6.docx'},
#                                 {
#                                     'searchapi_snippet': '3',
#                                     'searchapi_source': 'excel',
#                                     'searchapi_title': '幼儿园的保育员需要什么证（需要什么条件好不好做） -',
#                                     'searchapi_type': 'web',
#                                     'searchapi_url': 'https://renshichu.bit.edu.cn/docs//2026-03/65b2653627e7462a831cdf6c4bcf3325.xls'},
#                                 {
#                                     'searchapi_snippet': '4',
#                                     'searchapi_source': 'pdf',
#                                     'searchapi_title': '幼儿园的保育员需要什么证（需要什么条件好不好做） -',
#                                     'searchapi_type': 'web',
#                                     'searchapi_url': 'https://renshichu.bit.edu.cn/docs//2026-02/22cfe9cb32c240f78e9e7c2cfc4cde15.pdf'}
#                             ]}],
#       'institution_source_data': [],
#       'run_mode': 'Tuoyu'})

# # --- 4. 统一调度中心 (已重命名和扩展) ---
# class DataOrchestrator:
#     def __init__(self):
#         self.content_scrapers: Dict[str, ContentScraper] = {
#             "searchapi": SearchApiScraper(), "firecrawl": FirecrawlScraper(),
#             "jina": JinaScraper(), "tavily": TavilyScraper(),
#         }
#         self.job_scraper = ZhiLianJobScraper()
#         self.enterprise_scraper = TianyanEnterpriseScraper()
#
#     # async def process_all(self, url_list: List[Dict[str, str]], career_payload: Dict) -> Dict[str, Any]:
#     # ssl_context = httpx.create_ssl_context(verify=False)
#     # async with httpx.AsyncClient(http2=True, verify=ssl_context, timeout=30, follow_redirects=True,
#     #                              limits=httpx.Limits(max_connections=50)) as client:
#     #     # 创建两组任务
#     #     content_tasks = []
#     #     for item in url_list:
#     #         scraper = self.content_scrapers.get(item.get("provider"))
#     #         if scraper: content_tasks.append(scraper.scrape(item["url"], item["title"], client))
#
#     #     job_task = self.job_scraper.scrape_jobs(career_payload, client)
#
#     #     # 并发执行所有任务
#     #     results = await asyncio.gather(*content_tasks, job_task, return_exceptions=True)
#
#     #     # 分离结果
#     #     content_results = results[:-1]
#     #     job_result = results[-1]
#
#     #     return {"content_results": content_results, "job_result": job_result}
#     async def process_all(self, web_url_info_list: List[Dict[str, Any]], career_payload: Dict, enterprise_name: str) -> \
#             Dict[str, Any]:
#         ssl_context = httpx.create_ssl_context(verify=False)
#         async with httpx.AsyncClient(http2=True, verify=ssl_context, timeout=30, follow_redirects=True,
#                                      limits=httpx.Limits(max_connections=50)) as client:
#             content_tasks = []
#
#             # 【修改】从 web_url_info_list 创建抓取任务
#             for item in web_url_info_list:
#                 scraper = self.content_scrapers.get(item.get("provider"))
#                 # 默认使用 SearchApiScraper 作为备选
#                 if not scraper: scraper = self.content_scrapers["searchapi"]
#                 content_tasks.append(scraper.scrape(item, client))
#
#             job_task = self.job_scraper.scrape_jobs(career_payload, client)
#             enterprise_task = self.enterprise_scraper.scrape_enterprise(enterprise_name, client)
#
#             all_tasks = content_tasks + [job_task, enterprise_task]
#             results = await asyncio.gather(*all_tasks, return_exceptions=True)
#
#             content_results = results[:len(content_tasks)]
#             job_result = results[len(content_tasks)]
#             enterprise_result = results[len(content_tasks) + 1]
#             return {"content_results": content_results, "job_result": job_result,
#                     "enterprise_result": enterprise_result}
#
#
# # --- 5. Dify 节点主入口 ---
# async def main_async(raw_input: Any) -> Dict[str, Any]:
#     # 1. 解析输入
#     parsed_data = _parse_input_data(raw_input)
#     # 【修改】获取分离后的网页和视频列表
#     web_url_info_list = parsed_data["web_url_info_list"]
#     video_url_info_list = parsed_data["video_url_info_list"]
#     career_payload = parsed_data["career_payload"]
#     enterprise_name = parsed_data["enterprise_name"]
#
#     if not web_url_info_list and not video_url_info_list and not career_payload.get("keywords") and not enterprise_name:
#         print("🟡 所有输入均为空，提前返回。")
#         return {"scraped_datas": {}, "scraped_datas_str": "{}"}
#
#     # 2. 运行调度器 (只抓取网页内容)
#     orchestrator = DataOrchestrator()
#     results = await orchestrator.process_all(web_url_info_list, career_payload, enterprise_name)
#
#     # 3. 【修改】格式化网页内容输出，构建 all_source_list
#     all_source_list = []
#     for result in results["content_results"]:
#         if isinstance(result, Exception): continue
#         if result.get("status") == "success":  # 即使 content 为空也保留，以便下游判断
#             sanitized_url = re.sub(r'[^a-zA-Z0-9]', '-',
#                                    result.get("url", "").replace("https://", "").replace("http://", ""))
#             all_source_list.append({
#                 "type": "web",
#                 "source_id": f"web-{sanitized_url[:100]}",
#                 "url": result.get("url"),
#                 "title": result.get("title"),
#                 "source": result.get("source"),
#                 "snippet": result.get("snippet"),
#                 "query": result.get("query"),
#                 "content": result.get("content", "")  # 确保有 content 字段
#             })
#
#     # 【修改】格式化视频内容输出，构建 all_video_list
#     all_video_list = []
#     for video_item in video_url_info_list:
#         all_video_list.append({
#             "type": "video",
#             "url": video_item.get("url"),
#             "title": video_item.get("title"),
#             "source": video_item.get("source"),
#             "snippet": video_item.get("snippet"),
#             "video_id": video_item.get("video_id"),
#             "embed_url": video_item.get("embed_url"),
#             "thumbnail_url": video_item.get("thumbnail_url"),
#             "query": video_item.get("query")
#         })
#
#     # 4. 格式化招聘信息输出
#     career_postings = results["job_result"]
#     if isinstance(career_postings, Exception): career_postings = {"status": "failed", "data": [],
#                                                                   "message": f"任务异常: {career_postings}"}
#
#     # 5. 格式化企业信息输出
#     enterprise_info = results["enterprise_result"]
#     if isinstance(enterprise_info, Exception): enterprise_info = {"status": "failed", "data": None,
#                                                                   "message": f"任务异常: {enterprise_info}"}
#
#     # 6. 【修改】组装最终输出以符合新的数据结构
#     comprehensive_data_output = {
#         "all_source_list": all_source_list,
#         "all_video_list": all_video_list
#     }
#
#     final_output = {
#         "scraped_datas": {
#             "comprehensive_data": comprehensive_data_output,  # 修改此处的键和值
#             "career_postings": career_postings,
#             "enterprise_info": enterprise_info
#         }
#     }
#
#     return {
#         "scraped_datas": final_output["scraped_datas"],
#         "scraped_datas_str": json.dumps(final_output, ensure_ascii=False, indent=2)
#     }
#
#
# def main(datas_input: Any) -> Dict[str, Any]:
#     try:
#         return asyncio.run(main_async(raw_input=datas_input))
#     except Exception as e:
#         print(f"‼️ 节点执行时发生顶层错误: {e}")
#         # 【修改】错误负载以匹配新的 comprehensive_data 结构
#         error_payload = {
#             "comprehensive_data": {
#                 "all_source_list": [
#                     {"type": "web", "source_id": "NODE_EXECUTION_ERROR", "title": "节点执行失败", "url": "",
#                      "content": f"An error occurred: {str(e)}\n\n{traceback.format_exc()}"}],
#                 "all_video_list": []
#             },
#             "career_postings": {"status": "failed", "message": "节点执行失败", "data": []},
#             "enterprise_info": {"status": "failed", "message": "节点执行失败", "data": None}
#         }
#         return {
#             "scraped_datas": error_payload,
#             "scraped_datas_str": json.dumps({"scraped_datas": error_payload}, ensure_ascii=False, indent=2)
#         }

# async def main_async(raw_input: Any) -> Dict[str, Any]:
#     # 1. 解析输入
#     parsed_data = _parse_input_data(raw_input)
#     url_list = parsed_data["url_list"]
#     career_payload = parsed_data["career_payload"]
#     if not url_list and not career_payload.get("keywords"):
#         print("🟡 输入中没有有效的URL或招聘查询，提前返回。")
#         return {"scraped_datas": {}, "scraped_datas_str": "{}"}
#     enterprise_name = parsed_data["enterprise_name"]
#     # 2. 运行调度器
#     orchestrator = DataOrchestrator()
#     results = await orchestrator.process_all(url_list, career_payload, enterprise_name)

#     # 3. 格式化网页内容输出
#     comprehensive_content = []
#     for result in results["content_results"]:
#         if isinstance(result, Exception): continue
#         if result.get("status") == "success" and result.get("content"):
#             sanitized_url = re.sub(r'[^a-zA-Z0-9]', '-', result["url"].replace("https://", "").replace("http://", ""))
#             comprehensive_content.append({
#                 "source_id": f"web-{sanitized_url[:100]}", "source_name": result["title"],
#                 "url": result["url"], "content": result["content"]
#             })
#     # 4. 格式化招聘信息输出
#     career_postings = results["job_result"]
#     if isinstance(career_postings, Exception):
#         career_postings = {"status": "failed", "data": [], "message": f"任务异常: {career_postings}"}
#     # 5. 组装最终输出
#     final_output = {
#         "scraped_datas": {
#             "comprehensive_content": comprehensive_content,
#             "career_postings": career_postings
#         }
#     }
#     return {
#         "scraped_datas": final_output["scraped_datas"],
#         "scraped_datas_str": json.dumps(final_output, ensure_ascii=False, indent=2)
#     }

# def main(datas_input: Any) -> Dict[str, Any]:
#     try:
#         return asyncio.run(main_async(raw_input=datas_input))
#     except Exception as e:
#         print(f"‼️ 节点执行时发生顶层错误: {e}")
#         error_payload = {
#             "comprehensive_content": [{
#                 "source_id": "NODE_EXECUTION_ERROR", "source_name": "节点执行失败", "url": "",
#                 "content": f"An error occurred: {str(e)}\n\n{traceback.format_exc()}"
#             }],
#             "career_postings": {"status": "failed", "message": "节点执行失败", "data": []}
#         }
#         return {
#             "scraped_datas": error_payload,
#             "scraped_datas_str": json.dumps({"scraped_datas": error_payload}, ensure_ascii=False, indent=2)

# # Dify 依赖管理: 请确保已添加 httpx, json-repair, trafilatura, pypdf2, beautifulsoup4, lxml
# import asyncio
# import httpx
# import re
# import os
# import json
# import time
# import traceback
# from typing import Any, Dict, List, Literal, Optional
# from abc import ABC, abstractmethod
# from io import BytesIO
# from urllib.parse import urljoin

# # --- 核心依赖 ---
# # trafilatura 用于从HTML提取主要内容
# import trafilatura
# from trafilatura.settings import use_config

# # PyPDF2 用于解析PDF
# from PyPDF2 import PdfReader

# # BeautifulSoup 用于辅助解析HTML（例如提取视频）
# from bs4 import BeautifulSoup

# # --- 1. 输入解析模块 ---
# # 【已修复】替换这个函数
# def _parse_input_data(raw_input: Any) -> Dict[str, Any]:
#     """
#     健壮地解析上一个节点的输出，能同时处理带 "datas" 包装和不带包装的两种结构。
#     """
#     print(f"============== 步骤 1: 接收到原始输入 ==============\nTYPE: {type(raw_input)}\nVALUE: {raw_input}\n=======================================================")
#     if isinstance(raw_input, str):
#         if not raw_input.strip(): return {"url_list": [], "career_payload": {}, "enterprise_name": ""}
#         try:
#             data = json.loads(raw_input)
#         except json.JSONDecodeError as e:
#             raise ValueError(f"无法将输入字符串解析为JSON: {e}")
#     elif isinstance(raw_input, dict):
#         data = raw_input
#     else:
#         raise TypeError(f"期望的输入类型是 str 或 dict, 但收到了 {type(raw_input).__name__}")

#     # --- 核心修复逻辑 ---
#     # 检查顶层是否有 "datas" 键，如果没有，就认为当前整个对象就是我们要的数据体。
#     if "datas" in data and isinstance(data["datas"], dict):
#         print("  [解析器] 检测到 'datas' 包装层，将使用其内部数据。")
#         datas_obj = data["datas"]
#     else:
#         print("  [解析器] 未检测到 'datas' 包装层，将直接使用顶层数据。")
#         datas_obj = data
#     # --- 修复结束 ---

#     if not isinstance(datas_obj, dict): datas_obj = {}

#     comprehensive_data = datas_obj.get("comprehensive_data", [])
#     career_data = datas_obj.get("career_data", {})
#     tianyan_data = datas_obj.get("tianyan_check_data", "")

#     url_list = []
#     if isinstance(comprehensive_data, list):
#         for query_result in comprehensive_data:
#             if not isinstance(query_result, dict): continue
#             for res_list_key in ["web_results", "video_results"]:
#                 for res in query_result.get(res_list_key, []):
#                     if not isinstance(res, dict): continue
#                     url, title, provider = None, None, None
#                     for key, value in res.items():
#                         if key.endswith("_url"):
#                             url, provider = value, key.split('_url')[0]
#                             title = res.get(f"{provider}_title", "Untitled")
#                             break
#                     if url and provider:
#                         url_list.append({"url": url, "title": title, "provider": provider})
#     career_payload = career_data if isinstance(career_data, dict) else {}
#     enterprise_name = tianyan_data if isinstance(tianyan_data, str) else ""
#     parsed_result = {"url_list": url_list, "career_payload": career_payload, "enterprise_name": enterprise_name.strip()}

#     print(f"============== 步骤 2: 输入解析完毕 ==============\nURL 数量: {len(url_list)}\n招聘负载: {career_payload}\n企业名称: '{enterprise_name.strip()}'\n=======================================================")

#     return parsed_result

# # def _parse_input_data(raw_input: Any) -> Dict[str, Any]:
# #     """
# #     健壮地解析上一个节点的输出，分离出web搜索URL和招聘查询参数。
# #     """
# #     if isinstance(raw_input, str):
# #         if not raw_input.strip(): return {"url_list": [], "career_payload": {}}
# #         try:
# #             data = json.loads(raw_input)
# #         except json.JSONDecodeError as e:
# #             raise ValueError(f"无法将输入字符串解析为JSON: {e}")
# #     elif isinstance(raw_input, dict):
# #         data = raw_input
# #     else:
# #         raise TypeError(f"期望的输入类型是 str 或 dict, 但收到了 {type(raw_input).__name__}")
# #     # 安全地深入到 'datas' 结构
# #     datas_obj = data.get("datas", {})
# #     if not isinstance(datas_obj, dict): datas_obj = {}
# #     comprehensive_data = datas_obj.get("comprehensive_data", [])
# #     career_data = datas_obj.get("career_data", {})
# #     tianyan_data = datas_obj.get("tianyan_check_data", "")
# #     # 1. 提取URL列表
# #     url_list = []
# #     if isinstance(comprehensive_data, list):
# #         for query_result in comprehensive_data:
# #             if not isinstance(query_result, dict): continue
# #             for res_list_key in ["web_results", "video_results"]:
# #                 for res in query_result.get(res_list_key, []):
# #                     if not isinstance(res, dict): continue
# #                     url, title, provider = None, None, None
# #                     for key, value in res.items():
# #                         if key.endswith("_url"):
# #                             url = value
# #                             provider = key.split('_url')[0]
# #                             title = res.get(f"{provider}_title", "Untitled")
# #                             break
# #                     if url and provider:
# #                         url_list.append({"url": url, "title": title, "provider": provider})
# #     # 2. 提取职业查询负载
# #     career_payload = career_data if isinstance(career_data, dict) else {}

# #     # 3. 提取企业名称
# #     enterprise_name = tianyan_data if isinstance(tianyan_data, str) else ""
# #     return {"url_list": url_list, "career_payload": career_payload, "enterprise_name": enterprise_name.strip()}

# # --- 2. 抽象与实现分离：内容抓取器 ---
# class ContentScraper(ABC):
#     """内容抓取器的抽象基类。"""

#     @abstractmethod
#     async def scrape(self, url: str, title: str, client: httpx.AsyncClient) -> Dict[str, Any]:
#         """
#         抓取单个URL的内容。
#         成功时返回: {"url": str, "title": str, "content": str, "status": "success"}
#         失败时返回: {"url": str, "title": str, "content": "", "status": "failed", "error_message": str}
#         """
#         pass

# # --- 2.1 SearchAPI.io 的手动抓取与清洗实现 ---
# class SearchApiScraper(ContentScraper):
#     def __init__(self):
#         self.trafilatura_config = use_config()
#         self.trafilatura_config.set("DEFAULT", "EXTRACTION_TIMEOUT", "10")

#         # 编译常用的正则表达式以提高性能
#         self.NOISY_PATTERNS = [re.compile(p, re.IGNORECASE) for p in [
#             r'^\s*$', r'^[\-=*#_]{3,}$', r'.*\.(html|shtml|htm|php)\s*$',
#             r'.{0,50}(搜狐|网易|腾讯|新浪|登录|注册|版权所有|版权声明).{0,50}$',
#             r'\[\d+\]|\[下一页\]|\[上一页\]', r'\[(编辑|查看历史|讨论|阅读|来源|原标题)\]',
#             r'^\*+\s*\[.*?\]\(.*?\)',
#             r'^\s*(分享到|扫描二维码|返回搜狐|查看更多|责任编辑|记者|通讯员)',
#             r'^\s*([京公网安备京网文京ICP备]|互联网新闻信息服务许可证|信息网络传播视听节目许可证)',
#         ]]
#         self.IMG_PATTERN = re.compile(r'(!\[(.*?)\]\((.*?)\))')
#         self.LINK_PATTERN = re.compile(r'\[.*?\]\(.*?\)')
#         self.EDITOR_PATTERN = re.compile(r'(\(|\[)\s*责任编辑：.*?\s*(\)|\])')

#     # --- 2.1.1 内容提取工具 (来自您的代码) ---
#     def _extract_pdf_text(self, binary_content: bytes) -> str:
#         try:
#             reader = PdfReader(BytesIO(binary_content))
#             return "\n".join(page.extract_text() or "" for page in reader.pages)
#         except Exception as e:
#             print(f"⚠️ PDF 解析失败: {e}")
#             return ""

#     def _parse_videos_from_html(self, html: str, base_url: str) -> List[str]:
#         try:
#             soup = BeautifulSoup(html, "lxml")
#             videos = []
#             for video in soup.find_all("video"):
#                 src = video.get("src")
#                 if src: videos.append(urljoin(base_url, src))
#                 for source in video.find_all("source"):
#                     src = source.get("src")
#                     if src: videos.append(urljoin(base_url, src))
#             for iframe in soup.find_all("iframe"):
#                 src = iframe.get("src")
#                 if src and any(k in src for k in ["youtube", "vimeo", "embed", ".mp4"]):
#                     videos.append(urljoin(base_url, src))
#             return list(dict.fromkeys(videos))  # 去重并保持顺序
#         except Exception as e:
#             print(f"⚠️ 视频解析失败: {e}")
#             return []

#     # --- 2.1.2 内容清洗工具 (来自您的代码，已优化和异步化) ---
#     async def _is_valid_image_url_async(self, url: str, client: httpx.AsyncClient) -> bool:
#         if not url or not url.startswith(('http://', 'https://')): return False
#         try:
#             resp = await client.head(url, timeout=5, follow_redirects=True)
#             content_type = resp.headers.get('content-type', '').lower()
#             return resp.is_success and 'image' in content_type
#         except httpx.RequestError:
#             return False

#     async def _remove_invalid_images_async(self, md: str, client: httpx.AsyncClient) -> str:
#         MAX_IMAGES_TO_VALIDATE = 25
#         matches = list(self.IMG_PATTERN.finditer(md))
#         urls_to_check_all = {m.group(3).strip() for m in matches}

#         urls_to_check = set(list(urls_to_check_all)[:MAX_IMAGES_TO_VALIDATE])
#         if len(urls_to_check_all) > MAX_IMAGES_TO_VALIDATE:
#             print(f"⚠️ 图片数量过多 ({len(urls_to_check_all)}), 只验证前 {MAX_IMAGES_TO_VALIDATE} 张。")

#         tasks = {url: self._is_valid_image_url_async(url, client) for url in urls_to_check}
#         results = await asyncio.gather(*tasks.values(), return_exceptions=True)

#         url_status = dict(zip(tasks.keys(), results))
#         valid_urls = {url for url, res in url_status.items() if isinstance(res, bool) and res}
#         valid_urls.update(urls_to_check_all - urls_to_check)  # 未检查的默认有效

#         def replacer(match: re.Match):
#             return match.group(0) if match.group(3).strip() in valid_urls else ""

#         return self.IMG_PATTERN.sub(replacer, md)

#     def _is_noisy_line(self, line: str) -> bool:
#         stripped = line.strip()
#         for pat in self.NOISY_PATTERNS:
#             if pat.search(stripped): return True
#         links = self.LINK_PATTERN.findall(stripped)
#         if len(links) > 2 and len(stripped) / (len(links) + 1) < 30: return True
#         return False

#     async def _clean_content_async(self, text: str, client: httpx.AsyncClient) -> str:
#         if not text: return ""
#         text = await self._remove_invalid_images_async(text, client)

#         lines = text.splitlines()
#         cleaned_lines = []
#         for line in lines:
#             if not self._is_noisy_line(line):
#                 line = self.EDITOR_PATTERN.sub('', line).strip()
#                 if line: cleaned_lines.append(line)

#         # 去除连续空行
#         out = []
#         for i, line in enumerate(cleaned_lines):
#             if i > 0 and not line.strip() and not cleaned_lines[i - 1].strip():
#                 continue
#             out.append(line)

#         return "\n".join(out).strip()

#     # --- 2.1.3 主抓取函数 (来自您的代码，封装为scrape方法) ---
#     async def scrape(self, url: str, title: str, client: httpx.AsyncClient) -> dict:
#         print(f"🕸️ [SearchAPI Scraper] 开始处理: {url}")
#         try:
#             headers = {
#                 'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36'}
#             response = await client.get(url, timeout=20, headers=headers, follow_redirects=True)
#             response.raise_for_status()

#             final_url = str(response.url)
#             content_type = response.headers.get('content-type', '').lower()

#             raw_content, html_for_video_parsing = "", ""

#             if 'pdf' in content_type or final_url.lower().endswith(".pdf"):
#                 print(f"  📄 [SearchAPI Scraper] 检测到 PDF: {final_url}")
#                 pdf_bytes = await response.aread()
#                 raw_content = await asyncio.to_thread(self._extract_pdf_text, pdf_bytes)
#             else:
#                 print(f"  📑 [SearchAPI Scraper] 检测到 HTML: {final_url}")
#                 html_content = response.text
#                 html_for_video_parsing = html_content
#                 raw_content = await asyncio.to_thread(
#                     trafilatura.extract, html_content, config=self.trafilatura_config,
#                     output_format='markdown', include_images=True, favor_recall=True)

#             if not raw_content: raise ValueError("trafilatura 内容提取返回为空。")

#             print(f"  🧹 [SearchAPI Scraper] 正在清洗内容: {final_url}")
#             cleaned_content = await self._clean_content_async(raw_content, client)

#             if html_for_video_parsing:
#                 videos = self._parse_videos_from_html(html_for_video_parsing, final_url)
#                 if videos:
#                     video_section = "\n\n## 参考视频:\n" + "\n".join(f"- {vid}" for vid in videos)
#                     cleaned_content += video_section

#             print(f"✅ [SearchAPI Scraper] 成功: {url}")
#             return {"url": final_url, "title": title, "content": cleaned_content, "status": "success"}

#         except Exception as e:
#             error_msg = f"处理失败 {url}: {type(e).__name__} - {e}"
#             print(f"⚠️ [SearchAPI Scraper] {error_msg}")
#             return {"url": url, "title": title, "content": "", "status": "failed", "error_message": str(e)}

# # --- 2.2 FirecrawlScraper ---
# class FirecrawlScraper(ContentScraper):
#     def __init__(self):
#         self.api_key = os.environ.get("FIRECRAWL_API_KEY", "fc-a36b7d2fb273485680d0fe6abd686935")
#         if not self.api_key: raise ValueError("未提供 Firecrawl API Key。")
#         self.base_url = "https://api.firecrawl.dev/v2/scrape"
#         self.headers = {"Authorization": f"Bearer {self.api_key}", "Content-Type": "application/json"}

#     async def scrape(self, url: str, title: str, client: httpx.AsyncClient) -> dict:
#         print(f"🔥 [Firecrawl Scraper] 开始处理: {url}")
#         try:
#             # **关键修正：将 pageOptions 展平为顶级字段**
#             payload = {
#                 "url": url,
#                 "onlyMainContent": True,
#                 # 也可以在这里添加其他顶级选项，例如：
#                 "removeBase64Images": True,
#                 "blockAds": True
#             }
#             resp = await client.post(self.base_url, headers=self.headers, json=payload, timeout=45)

#             # 增加对4xx/5xx错误的详细日志记录
#             if not resp.is_success:
#                 try:
#                     error_details = resp.json()
#                     raise httpx.HTTPStatusError(f"API返回错误: {error_details}", request=resp.request, response=resp)
#                 except json.JSONDecodeError:
#                     resp.raise_for_status()  # 如果无法解析json，则抛出原始错误

#             data = resp.json()

#             # Firecrawl v2 的成功响应中没有 "success" 键，直接检查 data 字段
#             content_data = data.get("data", {})
#             if content_data is None:  # 可能是 null
#                 raise ValueError("API返回的 'data' 字段为 null。")

#             content = content_data.get("markdown")  # markdown 可能为空字符串，这是正常的
#             if content is None:
#                 raise ValueError("API未返回 'markdown' 字段。")

#             print(f"✅ [Firecrawl Scraper] 成功: {url}")
#             return {"url": url, "title": title, "content": content, "status": "success"}
#         except Exception as e:
#             error_msg = f"处理失败 {url}: {type(e).__name__} - {e}"
#             print(f"⚠️ [Firecrawl Scraper] {error_msg}")
#             return {"url": url, "title": title, "content": "", "status": "failed", "error_message": str(e)}

# # --- 2.3 JinaScraper ---
# class JinaScraper(ContentScraper):
#     def __init__(self):
#         self.api_key = os.environ.get("JINA_API_KEY",
#                                       "jina_b4348ffc39ca47bfbe753b95f59428c7i6ifkOFXRPdF3dRa5Rwb6T8FvrLH")
#         if not self.api_key: raise ValueError("未提供 Jina API Key。")
#         self.base_url = "https://r.jina.ai/"
#         self.headers = {
#             "Authorization": f"Bearer {self.api_key}",
#             "Content-Type": "application/json",
#             "Accept": "application/json",
#             "X-Return-Format": "markdown"  # 关键：直接获取Markdown
#         }

#     async def scrape(self, url: str, title: str, client: httpx.AsyncClient) -> dict:
#         print(f"🌀 [Jina Scraper] 开始处理: {url}")
#         try:
#             # Jina的Reader API有时对普通的GET请求更友好
#             target_url = f"{self.base_url}{url}"
#             resp = await client.get(target_url, headers=self.headers, timeout=45)
#             resp.raise_for_status()
#             content = resp.text
#             if not content: raise ValueError("API 返回内容为空。")

#             print(f"✅ [Jina Scraper] 成功: {url}")
#             return {"url": url, "title": title, "content": content, "status": "success"}

#         except Exception as e:
#             error_msg = f"处理失败 {url}: {type(e).__name__} - {e}"
#             print(f"⚠️ [Jina Scraper] {error_msg}")
#             return {"url": url, "title": title, "content": "", "status": "failed", "error_message": str(e)}

# # --- 2.4 TavilyScraper ---
# class TavilyScraper(ContentScraper):
#     def __init__(self):
#         self.api_key = os.environ.get("TAVILY_API_KEY", "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")
#         if not self.api_key: raise ValueError("未提供 Tavily API Key。")
#         self.base_url = "https://api.tavily.com/extract"
#         self.headers = {"Authorization": f"Bearer {self.api_key}", "Content-Type": "application/json"}

#     async def scrape(self, url: str, title: str, client: httpx.AsyncClient) -> dict:
#         print(f"🤖 [Tavily Scraper] 开始处理: {url}")
#         try:
#             # 注意：Tavily 的Python库没有extract方法，必须直接调用API
#             payload = {"urls": [url], "format": "markdown"}
#             resp = await client.post(self.base_url, json=payload, headers=self.headers, timeout=45)
#             resp.raise_for_status()
#             data = resp.json()

#             if not data.get("results") or not isinstance(data["results"], list):
#                 failed_info = data.get("failed_results", [])
#                 raise ValueError(f"API调用失败: {failed_info}")

#             result = data["results"][0]
#             content = result.get("raw_content")  # 文档显示raw_content，并且format=markdown
#             if not content: raise ValueError("API未返回raw_content内容。")

#             print(f"✅ [Tavily Scraper] 成功: {url}")
#             return {"url": url, "title": title, "content": content, "status": "success"}

#         except Exception as e:
#             error_msg = f"处理失败 {url}: {type(e).__name__} - {e}"
#             print(f"⚠️ [Tavily Scraper] {error_msg}")
#             return {"url": url, "title": title, "content": "", "status": "failed", "error_message": str(e)}

# # --- 2.5 ZhiLianJobScraper ---
# class ZhiLianJobScraper:
#     def __init__(self):
#         self.api_url = "http://119.45.167.133:12906/api/scrape/zhilian"
#         self.headers = {'accept': 'application/json', 'Content-Type': 'application/json'}

#     async def scrape_jobs(self, payload: Dict[str, Any], client: httpx.AsyncClient) -> Dict[str, Any]:
#         print(f"💼 [ZhiLian Scraper] 开始使用负载调用API: {json.dumps(payload, ensure_ascii=False)}")
#         if not payload or not payload.get("keywords") or not payload.get("provinces"):
#             msg = "负载无效，缺少 'keywords' 或 'provinces'。"
#             print(f"⚠️ [ZhiLian Scraper] {msg}")
#             return {"status": "skipped", "data": [], "message": msg}
#         try:
#             # 确保 page_size 是整数
#             if 'page_size' in payload: payload['page_size'] = int(payload['page_size'])

#             resp = await client.post(self.api_url, headers=self.headers, json=payload, timeout=60)
#             resp.raise_for_status()
#             response_data = resp.json()
#             if response_data.get("code") == 200:
#                 print(f"✅ [ZhiLian Scraper] 成功: {response_data.get('message')}")
#                 return {"status": "success", "data": response_data.get("data", []),
#                         "message": response_data.get("message")}
#             else:
#                 msg = f"API返回错误码 {response_data.get('code')}: {response_data.get('message')}"
#                 print(f"API returned non-200 code: {msg}")
#                 return {"status": "failed", "data": [], "message": msg}
#         except Exception as e:
#             error_msg = f"API请求失败: {type(e).__name__} - {e}"
#             print(f"⚠️ [ZhiLian Scraper] {error_msg}")
#             return {"status": "failed", "data": [], "message": error_msg}

# class TianyanEnterpriseScraper:
#     def __init__(self):
#         self.api_url = "http://open.api.tianyancha.com/services/open/ic/baseinfo/normal"
#         # 从环境变量或直接硬编码获取Token
#         self.token = os.environ.get("TIANYANCHA_TOKEN", "4d882100-ed23-4c22-a83b-c77af2e4be42")
#         self.headers = {'Authorization': self.token}
#     async def scrape_enterprise(self, name: str, client: httpx.AsyncClient) -> Dict[str, Any]:
#         """根据企业名称查询基本信息。"""
#         print(f"🏢 [Tianyan Scraper] 开始查询企业: {name}")
#         if not name:
#             msg = "企业名称为空，跳过查询。"
#             print(f"🟡 [Tianyan Scraper] {msg}")
#             return {"status": "skipped", "data": None, "message": msg}
#         try:
#             params = {"keyword": name}
#             resp = await client.get(self.api_url, headers=self.headers, params=params, timeout=30)
#             resp.raise_for_status()
#             response_data = resp.json()
#             if response_data.get("error_code") == 0:
#                 print(f"✅ [Tianyan Scraper] 成功查询到: {name}")
#                 return {"status": "success", "data": response_data.get("result"), "message": response_data.get("reason")}
#             else:
#                 msg = f"API返回错误码 {response_data.get('error_code')}: {response_data.get('reason')}"
#                 print(f"⚠️ [Tianyan Scraper] {msg}")
#                 return {"status": "failed", "data": None, "message": msg}
#         except Exception as e:
#             error_msg = f"API请求失败: {type(e).__name__} - {e}"
#             print(f"⚠️ [Tianyan Scraper] {error_msg}")
#             return {"status": "failed", "data": None, "message": error_msg}

# # --- 4. 统一调度中心 (已重命名和扩展) ---
# class DataOrchestrator:
#     def __init__(self):
#         self.content_scrapers: Dict[str, ContentScraper] = {
#             "searchapi": SearchApiScraper(), "firecrawl": FirecrawlScraper(),
#             "jina": JinaScraper(), "tavily": TavilyScraper(),
#         }
#         self.job_scraper = ZhiLianJobScraper()
#         self.enterprise_scraper = TianyanEnterpriseScraper()
#     # async def process_all(self, url_list: List[Dict[str, str]], career_payload: Dict) -> Dict[str, Any]:
#         # ssl_context = httpx.create_ssl_context(verify=False)
#         # async with httpx.AsyncClient(http2=True, verify=ssl_context, timeout=30, follow_redirects=True,
#         #                              limits=httpx.Limits(max_connections=50)) as client:
#         #     # 创建两组任务
#         #     content_tasks = []
#         #     for item in url_list:
#         #         scraper = self.content_scrapers.get(item.get("provider"))
#         #         if scraper: content_tasks.append(scraper.scrape(item["url"], item["title"], client))

#         #     job_task = self.job_scraper.scrape_jobs(career_payload, client)

#         #     # 并发执行所有任务
#         #     results = await asyncio.gather(*content_tasks, job_task, return_exceptions=True)

#         #     # 分离结果
#         #     content_results = results[:-1]
#         #     job_result = results[-1]

#         #     return {"content_results": content_results, "job_result": job_result}
#     async def process_all(self, url_list: List[Dict[str, str]], career_payload: Dict, enterprise_name: str) -> Dict[str, Any]:
#         ssl_context = httpx.create_ssl_context(verify=False)
#         async with httpx.AsyncClient(http2=True, verify=ssl_context, timeout=30, follow_redirects=True, limits=httpx.Limits(max_connections=50)) as client:
#             # 创建三组任务
#             content_tasks, job_task, enterprise_task = [], None, None

#             for item in url_list:
#                 scraper = self.content_scrapers.get(item.get("provider"))
#                 if scraper: content_tasks.append(scraper.scrape(item["url"], item["title"], client))

#             job_task = self.job_scraper.scrape_jobs(career_payload, client)
#             enterprise_task = self.enterprise_scraper.scrape_enterprise(enterprise_name, client) # 【新增】

#             # 并发执行所有任务
#             all_tasks = content_tasks + [job_task, enterprise_task]
#             results = await asyncio.gather(*all_tasks, return_exceptions=True)

#             # 分离结果
#             content_results = results[:len(content_tasks)]
#             job_result = results[len(content_tasks)]
#             enterprise_result = results[len(content_tasks) + 1] # 【新增】
#             return {"content_results": content_results, "job_result": job_result, "enterprise_result": enterprise_result}

# # --- 5. Dify 节点主入口 ---
# async def main_async(raw_input: Any) -> Dict[str, Any]:
#     # 1. 解析输入
#     parsed_data = _parse_input_data(raw_input)
#     url_list = parsed_data["url_list"]
#     career_payload = parsed_data["career_payload"]
#     enterprise_name = parsed_data["enterprise_name"] # 【新增】

#     if not url_list and not career_payload.get("keywords") and not enterprise_name:
#         print("🟡 所有输入均为空，提前返回。")
#         return {"scraped_datas": {}, "scraped_datas_str": "{}"}
#     # 2. 运行调度器
#     orchestrator = DataOrchestrator()
#     results = await orchestrator.process_all(url_list, career_payload, enterprise_name)
#     # 3. 格式化网页内容输出
#     comprehensive_content = []
#     for result in results["content_results"]:
#         if isinstance(result, Exception): continue
#         if result.get("status") == "success" and result.get("content"):
#             sanitized_url = re.sub(r'[^a-zA-Z0-9]', '-', result["url"].replace("https://", "").replace("http://", ""))
#             comprehensive_content.append({"source_id": f"web-{sanitized_url[:100]}", "source_name": result["title"], "url": result["url"], "content": result["content"]})
#     # 4. 格式化招聘信息输出
#     career_postings = results["job_result"]
#     if isinstance(career_postings, Exception): career_postings = {"status": "failed", "data": [], "message": f"任务异常: {career_postings}"}
#     # 5. 【新增】格式化企业信息输出
#     enterprise_info = results["enterprise_result"]
#     if isinstance(enterprise_info, Exception): enterprise_info = {"status": "failed", "data": None, "message": f"任务异常: {enterprise_info}"}
#     # 6. 【调整】组装最终输出
#     final_output = {
#         "scraped_datas": {
#             "comprehensive_content": comprehensive_content,
#             "career_postings": career_postings,
#             "enterprise_info": enterprise_info
#         }
#     }
#     return {
#         "scraped_datas": final_output["scraped_datas"],
#         "scraped_datas_str": json.dumps(final_output, ensure_ascii=False, indent=2)
#     }
# # 【调整】main 函数
# def main(datas_input: Any) -> Dict[str, Any]:
#     try:
#         return asyncio.run(main_async(raw_input=datas_input))
#     except Exception as e:
#         print(f"‼️ 节点执行时发生顶层错误: {e}")
#         error_payload = {
#             "comprehensive_content": [{"source_id": "NODE_EXECUTION_ERROR", "source_name": "节点执行失败", "url": "", "content": f"An error occurred: {str(e)}\n\n{traceback.format_exc()}"}],
#             "career_postings": {"status": "failed", "message": "节点执行失败", "data": []},
#             "enterprise_info": {"status": "failed", "message": "节点执行失败", "data": None}
#         }
#         return {
#             "scraped_datas": error_payload,
#             "scraped_datas_str": json.dumps({"scraped_datas": error_payload}, ensure_ascii=False, indent=2)
#         }

# # async def main_async(raw_input: Any) -> Dict[str, Any]:
# #     # 1. 解析输入
# #     parsed_data = _parse_input_data(raw_input)
# #     url_list = parsed_data["url_list"]
# #     career_payload = parsed_data["career_payload"]
# #     if not url_list and not career_payload.get("keywords"):
# #         print("🟡 输入中没有有效的URL或招聘查询，提前返回。")
# #         return {"scraped_datas": {}, "scraped_datas_str": "{}"}
# #     enterprise_name = parsed_data["enterprise_name"]
# #     # 2. 运行调度器
# #     orchestrator = DataOrchestrator()
# #     results = await orchestrator.process_all(url_list, career_payload, enterprise_name)

# #     # 3. 格式化网页内容输出
# #     comprehensive_content = []
# #     for result in results["content_results"]:
# #         if isinstance(result, Exception): continue
# #         if result.get("status") == "success" and result.get("content"):
# #             sanitized_url = re.sub(r'[^a-zA-Z0-9]', '-', result["url"].replace("https://", "").replace("http://", ""))
# #             comprehensive_content.append({
# #                 "source_id": f"web-{sanitized_url[:100]}", "source_name": result["title"],
# #                 "url": result["url"], "content": result["content"]
# #             })
# #     # 4. 格式化招聘信息输出
# #     career_postings = results["job_result"]
# #     if isinstance(career_postings, Exception):
# #         career_postings = {"status": "failed", "data": [], "message": f"任务异常: {career_postings}"}
# #     # 5. 组装最终输出
# #     final_output = {
# #         "scraped_datas": {
# #             "comprehensive_content": comprehensive_content,
# #             "career_postings": career_postings
# #         }
# #     }
# #     return {
# #         "scraped_datas": final_output["scraped_datas"],
# #         "scraped_datas_str": json.dumps(final_output, ensure_ascii=False, indent=2)
# #     }

# # def main(datas_input: Any) -> Dict[str, Any]:
# #     try:
# #         return asyncio.run(main_async(raw_input=datas_input))
# #     except Exception as e:
# #         print(f"‼️ 节点执行时发生顶层错误: {e}")
# #         error_payload = {
# #             "comprehensive_content": [{
# #                 "source_id": "NODE_EXECUTION_ERROR", "source_name": "节点执行失败", "url": "",
# #                 "content": f"An error occurred: {str(e)}\n\n{traceback.format_exc()}"
# #             }],
# #             "career_postings": {"status": "failed", "message": "节点执行失败", "data": []}
# #         }
# #         return {
# #             "scraped_datas": error_payload,
# #             "scraped_datas_str": json.dumps({"scraped_datas": error_payload}, ensure_ascii=False, indent=2)
